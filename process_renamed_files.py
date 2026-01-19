"""
Process the 4 renamed .ppt -> .pptx files
"""
import os
import sys
import json
import boto3
from pathlib import Path
from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'backend'))

from backend.bedrock_embeddings import BedrockEmbeddings

print("="*80)
print("PROCESSING RENAMED .PPT -> .PPTX FILES")
print("="*80)

# Config
S3_BUCKET = "smart-tutor-vector-bucket"
S3_REGION = "us-east-1"
CHUNK_SIZE = 512

# Initialize
print("\n1. Initializing Bedrock...")
sys.stdout.flush()
bedrock = BedrockEmbeddings()
print("   ✅ Bedrock ready")
sys.stdout.flush()

print("\n   Initializing S3 client...")
sys.stdout.flush()
s3 = boto3.client('s3', region_name=S3_REGION)
print("   ✅ S3 client ready")
sys.stdout.flush()

# Files to process
files_to_process = [
    "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules/Module 5/Lesson Five- Data cleaning and preprocessing.pptx",
    "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules/Module 12/INFO 5731 - Lesson nine - Sentiment analysis-1.pptx",
    "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules/Module 6/Lesson six-Feature extraction from text-2024.pptx",
    "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules/Module 8/Lesson Seven- Information Extraction from Textual Data_Updated-02262024 (1).pptx"
]

def extract_text_from_pptx(pptx_path):
    """Extract text from PPTX"""
    texts = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                texts.append("\n".join(slide_text))
    except Exception as e:
        print(f"      ✗ PPTX error: {e}")
    return "\n\n".join(texts)

def simple_chunk(text, size=512):
    """Simple chunking by characters with overlap"""
    chunks = []
    overlap = size // 4  # 25% overlap
    for i in range(0, len(text), size - overlap):
        chunk = text[i:i+size]
        if len(chunk.strip()) > 50:  # Only keep meaningful chunks
            chunks.append(chunk)
    return chunks

# Process files
print(f"\n2. Processing {len(files_to_process)} renamed files...")
sys.stdout.flush()
uploaded = 0
failed_files = 0

for idx, file_path in enumerate(files_to_process, 1):
    file = os.path.basename(file_path)
    print(f"\n   [{idx}/{len(files_to_process)}] Processing: {file}")
    sys.stdout.flush()

    if not os.path.exists(file_path):
        print(f"      ✗ File not found: {file_path}")
        failed_files += 1
        continue

    try:
        # Extract text
        text = extract_text_from_pptx(file_path)

        if not text or len(text.strip()) < 100:
            print(f"      → Skipped (no text extracted)")
            continue

        # Chunk
        chunks = simple_chunk(text, CHUNK_SIZE)
        print(f"      → {len(chunks)} chunks created")
        sys.stdout.flush()

        # Upload each chunk
        for i, chunk_text in enumerate(chunks):
            try:
                # Generate embedding
                embedding = bedrock._encode_single(chunk_text)

                # Create chunk
                chunk_id = f"{Path(file).stem}_chunk_{i:04d}"
                chunk_data = {
                    'chunk_id': chunk_id,
                    'text': chunk_text,
                    'embedding': embedding,
                    'source_file': file,
                    'chunk_index': i
                }

                # Upload to S3
                s3.put_object(
                    Bucket=S3_BUCKET,
                    Key=f"chunks/{chunk_id}.json",
                    Body=json.dumps(chunk_data),
                    ContentType='application/json'
                )

                uploaded += 1

                if uploaded % 10 == 0:
                    print(f"      → Progress: {uploaded} chunks uploaded...")
                    sys.stdout.flush()

            except Exception as e:
                print(f"      ✗ Chunk {i} error: {e}")

    except Exception as e:
        print(f"      ✗ File error: {e}")
        failed_files += 1

print(f"\n{'='*80}")
print(f"✅ COMPLETE!")
print(f"{'='*80}")
print(f"Uploaded: {uploaded} chunks")
print(f"Failed files: {failed_files}")
print(f"Location: s3://{S3_BUCKET}/chunks/")
print(f"\nTotal chunks in S3 now: 14177 + {uploaded} = {14177 + uploaded}")
print(f"\nNext: Rebuild S3 vector index and restart backend")
