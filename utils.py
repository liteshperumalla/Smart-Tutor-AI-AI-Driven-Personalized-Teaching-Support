try:
    import streamlit as st
except ImportError:
    st = None  # Running outside Streamlit (e.g. FastAPI backend)
import os
import re
import base64
import json
import textwrap
from pathlib import Path
import logging
import smtplib
from email.message import EmailMessage
import time
import datetime
from functools import lru_cache
from typing import Optional, List, Dict, Any, Union, Tuple
import requests
from urllib.parse import quote
from youtube_transcript_api import (
    YouTubeTranscriptApi,
    TranscriptsDisabled,
    NoTranscriptFound,
)
from bs4 import BeautifulSoup
from backend.services.status_service import (
    get_knowledge_base_stats,
    get_system_status,
    check_ollama_status,
)

# LlamaIndex imports
from llama_index.core import (
    StorageContext,
    load_index_from_storage,
    get_response_synthesizer,
    VectorStoreIndex,
    PromptTemplate,
    Settings,
)
from llama_index.core.schema import Document, TextNode, NodeWithScore

# S3 Vector Store imports
from backend.s3_retriever import create_s3_retriever
from backend.config import config
from backend.llm_provider import get_llm

# File processing imports
from pptx import Presentation
from docx import Document as DocxDoc
import mammoth
from fpdf import FPDF
from PIL import Image
import pytesseract

try:
    from serpapi import GoogleSearch
except ImportError:
    GoogleSearch = None

# Import langfuse for tracing
try:
    from langfuse import Langfuse

    langfuse_client = Langfuse()
except ImportError:
    logging.warning("Langfuse not available. Tracing will be disabled.")
    langfuse_client = None
except Exception as e:
    logging.error(f"Failed to initialize Langfuse client: {e}")
    langfuse_client = None

# --- Constants ---
PREV_CHAT_DIR = config.PREV_CHAT_DIR
os.makedirs(PREV_CHAT_DIR, exist_ok=True)

# Configure logging
logging.basicConfig(
    level=logging.ERROR, format="%(asctime)s - %(levelname)s - %(message)s"
)

# --- LlamaIndex Setup ---
PERSIST_DIR = "./persisted_index"
os.makedirs(PERSIST_DIR, exist_ok=True)
CHROMA_DB_PATH = os.getenv("CHROMA_DB_PATH", "./chroma_db")

# Initialize LLM for response generation
try:
    Settings.llm = get_llm()
    logging.info(f"LLM initialized for response generation: {config.LLM_PROVIDER}")
except Exception as e:
    logging.error(f"Failed to initialize LLM: {e}")
    Settings.llm = None

WEB_SEARCH_ENABLED = os.getenv("WEB_SEARCH_ENABLED", "true").lower() == "true"
SERPAPI_API_KEY = os.getenv("SERPAPI_API_KEY", "").strip()
MAX_WEB_RESULTS = int(os.getenv("MAX_WEB_RESULTS", "3"))


CHAT_QA_TEMPLATE = PromptTemplate(
    (
        "You are Smart AI Tutor, a knowledgeable teaching assistant for graduate students. "
        "Using only the provided context, craft a brief yet thorough explanation that defines the idea, "
        "highlights why it matters, and includes a concrete example or analogy students can relate to. "
        "IMPORTANT: When citing information from the sources, include inline citations using the format [1], [2], etc. "
        "corresponding to the source numbers in the context below. Place the citation immediately after the relevant statement. "
        "Use **bold text** to highlight key terms and important concepts. "
        "Clearly state when information is missing. "
        "Never reply with a single number, bullet label, or one-word answer—always respond in at least "
        "three complete sentences.\n"
        "---------------------\n"
        "{context_str}\n"
        "---------------------\n"
        "Question: {query_str}\n"
        "Answer (remember to use [1], [2], etc. for inline citations and **bold** for key terms):"
    )
)

WEB_SEARCH_TEMPLATE = PromptTemplate(
    (
        "You are Smart AI Tutor, and you are providing information from web search results because "
        "the coursework materials did not contain the answer. "
        "Always open with the sentence: 'Information from web search (not found in course materials)' "
        "followed by a blank line. Summarize the results in clear academic language. "
        "IMPORTANT: Include inline citations using the format [1], [2], etc. corresponding to the source numbers below. "
        "Place the citation immediately after the relevant statement. "
        "Use **bold text** to highlight key terms and important concepts. "
        "Note when multiple sources agree or disagree. "
        "---------------------\n"
        "{context_str}\n"
        "---------------------\n"
        "Question: {query_str}\n"
        "Answer (remember to use [1], [2], etc. for inline citations and **bold** for key terms):"
    )
)


def _should_search_web(query: str, context: str, sources: List[Dict[str, Any]]) -> bool:
    """Return True when retrieved coursework context looks insufficient."""
    if not WEB_SEARCH_ENABLED:
        return False
    if not context or not context.strip():
        return True

    # If we got results from S3/local RAG, check token coverage (not exact phrase match)
    # This allows semantic search to work - S3 embeddings find relevant content
    # even if it doesn't contain the exact query phrase
    context_text = context.lower()
    tokens = [
        token for token in re.findall(r"[a-zA-Z0-9]+", query.lower()) if len(token) > 3
    ]
    if not tokens:
        return False

    hits = sum(1 for token in tokens if token in context_text)
    coverage = hits / len(tokens)

    # Only fall back to web search if coverage is very low (< 30%)
    # This trusts the semantic search results from S3/embeddings
    if coverage < 0.3:
        return True

    return False


def _search_with_serpapi(query: str, max_results: int):
    if not SERPAPI_API_KEY or GoogleSearch is None:
        return []
    try:
        search = GoogleSearch(
            {
                "q": query,
                "engine": "google",
                "api_key": SERPAPI_API_KEY,
                "num": max_results,
                "safe": "active",
                "hl": "en",
                "gl": "us",
            }
        )
        result = search.get_dict()
        formatted = []
        if "answer_box" in result:
            box = result["answer_box"]
            formatted.append(
                {
                    "title": box.get("title", "Featured Answer"),
                    "content": box.get("answer") or box.get("snippet", ""),
                    "url": box.get("link", ""),
                    "score": 1.0,
                    "source": "Google Featured Answer",
                }
            )
        for idx, entry in enumerate(result.get("organic_results", [])[:max_results]):
            formatted.append(
                {
                    "title": entry.get("title") or "Result",
                    "content": entry.get("snippet", ""),
                    "url": entry.get("link", ""),
                    "score": 0.9 - (idx * 0.1),
                    "source": entry.get("source", "Google Search"),
                }
            )
        return [res for res in formatted if res["content"]]
    except Exception as exc:
        logging.error(f"SerpAPI search failed: {exc}")
        return []


def _search_with_duckduckgo(query: str, max_results: int):
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/118.0 Safari/537.36"
            )
        }
        resp = requests.get(
            "https://html.duckduckgo.com/html/",
            params={"q": query},
            headers=headers,
            timeout=10,
        )
        resp.raise_for_status()
        soup = BeautifulSoup(resp.content, "html.parser")
        results = []
        for entry in soup.find_all("div", class_="result")[:max_results]:
            title_el = entry.find("a", class_="result__a")
            snippet_el = entry.find("a", class_="result__snippet")
            if not title_el or not snippet_el:
                continue
            results.append(
                {
                    "title": title_el.get_text(strip=True) or "Result",
                    "content": snippet_el.get_text(strip=True),
                    "url": title_el.get("href", ""),
                    "score": 0.7,
                    "source": "DuckDuckGo",
                }
            )
        return results
    except Exception as exc:
        logging.error(f"DuckDuckGo search failed: {exc}")
        return []


def search_web_results(query: str, max_results: int = MAX_WEB_RESULTS):
    """Fetch web search snippets using SerpAPI when available, otherwise DuckDuckGo."""
    if not WEB_SEARCH_ENABLED:
        return []
    results = _search_with_serpapi(query, max_results)
    if results:
        return results
    return _search_with_duckduckgo(query, max_results)


def get_storage_context():
    """Get storage context with proper error handling"""
    try:
        return StorageContext.from_defaults(persist_dir=PERSIST_DIR)
    except Exception as e:
        if st:
            st.error(f"Error initializing storage context: {e}")
        logging.error(f"Storage context initialization failed: {e}")
        return None


# --- CSS Definitions ---
LIGHT_MODE_CSS = """
<style>
  /* Headings & paragraphs black */
  h1, h2, h3, h4, p { color: black !important; }
  html, body, [class*="css"] { background-color: #f9fcff !important; color: black !important; }
  .main-title, .subtitle, .course-topics, .professor-title, .professor-name,
  .chat-title, .chat-subtitle { color: black !important; }
  .disclaimer {
    text-align: center;
    font-size: 0.8em;
    color: gray;
    margin-top: 20px; 
    padding-bottom: 10px;
  }
  /* Sidebar Button Styles */
  [data-testid="stSidebar"] button {
      background-color: transparent !important;
      color: inherit !important;
      border: none !important;
      text-align: left !important;
      padding: 10px 15px !important;
      border-radius: 8px !important;
      transition: background-color 0.2s;
  }
  [data-testid="stSidebar"] button:hover {
      background-color: #e0e0e0 !important;
      color: black !important;
  }
  [data-testid="stSidebar"] button:focus {
      background-color: #d6d6d6 !important;
  }
  [data-testid="stSidebar"] button:active {
      background-color: #cfcfcf !important;
  }
</style>
"""

DARK_MODE_CSS = """
<style>
  h1, h2, h3, h4, p { color: white !important; }
  [data-testid="stAppViewContainer"], [data-testid="stSidebar"],
  [data-testid="stHeader"], [data-testid="stToolbar"] {
      background-color: #000 !important;
      color: #fff !important;
  }
  .main-title, .subtitle, .course-topics,
  .professor-title, .professor-name,
  .chat-title, .chat-subtitle { color: #fff !important; }
  .disclaimer {
    text-align: center;
    font-size: 0.8em;
    color: gray;
    margin-top: 20px;
    padding-bottom: 10px;
  }
</style>
"""

BASE_CSS = """
<style>
  h1, h2, h3, h4, p { color: inherit !important; }
  .main-title { font-size: 3em; font-weight: 800; margin-bottom: 10px; }
  hr { border: none; height: 3px; background-color: #e0e0e0; margin-bottom: 25px; }
  .announcement-card { background: #fff3e0; border-left: 6px solid #fb8c00;
    border-radius: 8px; padding: 16px 20px; margin-bottom: 25px; max-height: 150px;
    overflow-y: auto;
  }
  .topics-card { background: #fff; border-radius: 8px; padding: 20px; margin-bottom: 25px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.08);
  }
  .professor-card { background: #fff; padding: 20px; border-radius: 8px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.1); text-align: center;
  }
  .professor-title { margin-top: 0; }
  .professor-name  { font-weight: 600; font-size: 1.1em; margin: 0; }
  .chat-container { display: inline-block; max-width: 70%; padding: 10px 15px;
    border-radius: 8px; margin-bottom: 10px; white-space: pre-wrap;
    word-wrap: break-word;
  }
  .user-chat      { background-color: #e3f2fd; margin-left: auto; }
  .assistant-chat { background-color: #f1f8e9; margin-right: auto; }
  .voice-button { margin-left: 10px; cursor: pointer; border: none;
    background-color: transparent; font-size: 20px;
  }
</style>
"""


# --- Helper Functions ---
def render_footer():
    """Render a footer with disclaimer"""
    if not st:
        return
    st.markdown(
        """
        <style>
        .footer {
            position: fixed;
            bottom: 0;
            left: 0;
            right: 0;
            background-color: #f1f1f1;
            padding: 10px 20px;
            text-align: center;
            font-size: 14px;
            color: #555;
            z-index: 999;
            border-top: 1px solid #ccc;
        }
        </style>
        <div class="footer">
            Disclaimer: The Smart AI Tutor may occasionally make mistakes. Please verify important information independently.
        </div>
        """,
        unsafe_allow_html=True,
    )


def _safe_name(name: str) -> str:
    """Convert name to safe filename"""
    return re.sub(r"[^a-zA-Z0-9_-]", "_", name)


def _chat_path(name: str) -> str:
    """Get full path for chat session file"""
    return os.path.join(PREV_CHAT_DIR, f"{_safe_name(name)}.json")


def resolve_chat_file_path(name: str) -> Path:
    """Find an existing chat file path using common naming variants."""
    candidates = [
        Path(PREV_CHAT_DIR) / f"{name}.json",
        Path(_chat_path(name)),
        Path(PREV_CHAT_DIR) / f"{sanitize_filename(name)}.json",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[1]


def load_chat_sessions() -> Dict[str, List]:
    """Load all chat sessions from disk"""
    sessions = {}

    if not os.path.exists(PREV_CHAT_DIR):
        os.makedirs(PREV_CHAT_DIR)
        return sessions

    for fname in os.listdir(PREV_CHAT_DIR):
        if not fname.endswith(".json"):
            continue
        path = os.path.join(PREV_CHAT_DIR, fname)
        try:
            with open(path, "r", encoding="utf-8") as f:
                session_name = fname[:-5]  # Remove .json extension
                sessions[session_name] = json.load(f)
        except Exception as e:
            logging.error(f"Error loading chat session {fname}: {e}")

    return sessions


def get_recent_chat_summaries(limit: int = 3) -> List[Dict[str, Any]]:
    """Return lightweight metadata for the most recent chat sessions."""
    sessions = load_chat_sessions()
    summaries: List[Dict[str, Any]] = []

    for name, history in sessions.items():
        preview_text = ""
        if history:
            for entry in reversed(history):
                if isinstance(entry, dict):
                    preview_text = entry.get("content", "") or ""
                elif isinstance(entry, (list, tuple)) and len(entry) >= 2:
                    preview_text = entry[1] or ""
                if preview_text:
                    preview_text = preview_text.strip()
                    if preview_text:
                        break

        file_path = resolve_chat_file_path(name)
        try:
            updated_at = file_path.stat().st_mtime
        except OSError:
            updated_at = 0

        if updated_at:
            updated_display = datetime.datetime.fromtimestamp(updated_at).strftime(
                "%b %d, %I:%M %p"
            )
        else:
            updated_display = "Not saved yet"

        summaries.append(
            {
                "name": name,
                "preview": textwrap.shorten(preview_text, width=110, placeholder="...")
                if preview_text
                else "No messages yet",
                "updated_at": updated_at,
                "updated_display": updated_display,
            }
        )

    summaries.sort(key=lambda item: item["updated_at"], reverse=True)
    if limit is None or limit <= 0:
        return summaries
    return summaries[:limit]


def save_chat_session(name: str, history: List) -> None:
    """Save chat session to disk"""
    formatted_history = []
    for entry in history:
        if isinstance(entry, (list, tuple)):
            if len(entry) == 3:  # user: (role, message, timestamp)
                role, text, timestamp_val = entry
                formatted_history.append(
                    {
                        "role": role,
                        "content": text,
                        "timestamp": timestamp_val,
                        "sources": [],
                    }
                )
            elif len(entry) == 4:  # assistant: (role, message, sources, timestamp)
                role, text, sources_val, timestamp_val = entry
                formatted_history.append(
                    {
                        "role": role,
                        "content": text,
                        "timestamp": timestamp_val,
                        "sources": sources_val,
                    }
                )
        elif isinstance(entry, dict):  # Already correct format
            formatted_history.append(entry)

    filepath = resolve_chat_file_path(name)
    try:
        with filepath.open("w", encoding="utf-8") as f:
            json.dump(formatted_history, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logging.error(f"Error saving chat session {name} to {filepath}: {e}")


def sanitize_filename(name: str) -> str:
    """Sanitize filename for safe storage"""
    safe = re.sub(r"[^0-9A-Za-z_-]", "_", name).lower()
    return safe[:50].rstrip("_")


def is_greeting(msg: str) -> bool:
    """Check if message is a greeting"""
    greetings = [
        "hi",
        "hello",
        "hey",
        "good morning",
        "good evening",
        "greetings",
        "thanks",
        "thank you",
        "bye",
        "goodbye",
    ]
    return any(msg.lower().strip().startswith(g) for g in greetings)


def generate_response_with_sources(
    query: str, user_id: str = "default-user", session_id: Optional[str] = None
) -> Tuple[str, List[Dict]]:
    """Generate response with sources (non-streaming version)"""
    main_trace = None
    if langfuse_client:
        main_trace = langfuse_client.trace(
            name="chat-rag-interaction",
            user_id=user_id,
            session_id=session_id,
            input={"query": query},
            tags=["chat", "RAG"],
        )

    response_sources = []

    try:
        if is_greeting(query):
            response = "Hello! How can I assist you today?"
            if main_trace:
                main_trace.update(output={"response": response})
            return response, []

        # Use S3 vectors or local ChromaDB
        if config.USE_S3_VECTORS:
            try:
                retriever = create_s3_retriever(similarity_top_k=3)
                retrieved_nodes = retriever.retrieve(query)
            except Exception as e:
                err_msg = f"Error using S3 vectors: {str(e)}"
                if main_trace:
                    main_trace.update(output={"error": err_msg}, level="ERROR")
                return err_msg, []
        else:
            storage_context = get_storage_context()
            if storage_context is None:
                err_msg = "Error: Storage context not initialized. Cannot query index."
                if main_trace:
                    main_trace.update(output={"error": err_msg}, level="ERROR")
                return err_msg, []

            idx = load_index_from_storage(storage_context)
            retriever = idx.as_retriever(similarity_top_k=3)
            retrieved_nodes = retriever.retrieve(query)

        for n_ws in retrieved_nodes:
            node = n_ws.node
            source_text = (
                node.get_text()
                if hasattr(node, "get_text")
                else (node.text if hasattr(node, "text") else "")
            )

            # Check both file_path (ChromaDB) and source_file (S3)
            fp = node.metadata.get("file_path") or node.metadata.get("source_file")
            file_name = os.path.basename(fp) if fp else "Unknown Source"

            # Generate S3 document URL instead of local file path
            source_url = (
                f"/api/backend/files/s3-document?source_file={quote(file_name)}"
                if fp
                else None
            )

            response_sources.append(
                {
                    "file_name": file_name,
                    "file_path": fp,  # Keep for backwards compatibility
                    "source_url": source_url,  # New S3 URL for frontend
                    "page": node.metadata.get("page_number"),
                    "slide": node.metadata.get("slide_number"),
                    "chunk_text": source_text[:300] + "..."
                    if len(source_text) > 300
                    else source_text,
                }
            )

        synth = get_response_synthesizer(response_mode="compact")
        response_obj = synth.synthesize(query=query, nodes=retrieved_nodes)
        response_text = str(response_obj)

        if main_trace:
            main_trace.update(
                output={"response": response_text},
                metadata={
                    "num_retrieved_sources": len(response_sources),
                    "retrieved_source_sample": [
                        s["file_name"] for s in response_sources[:2]
                    ],
                },
            )

        return response_text, response_sources

    except Exception as e:
        logging.error(f"Error in generate_response_with_sources: {e}", exc_info=True)
        if main_trace:
            main_trace.update(output={"error": str(e)}, level="ERROR")
        error_msg = f"⚠️ Error processing your query: {e}"
        return error_msg, response_sources


def generate_response_stream_and_sources(
    query: str,
    user_id: str = "default-user",
    session_id: Optional[str] = None,
    enable_llm_judge: bool = False,
    model_id: Optional[str] = None,
):
    """Generate streaming response with sources

    Args:
        query: The user's question
        user_id: User identifier for tracing
        session_id: Session identifier for tracing
        enable_llm_judge: Whether to enable LLM judge evaluation
        model_id: Optional AWS Bedrock model ID to use (e.g., 'anthropic.claude-3-5-sonnet-20241022-v2:0')
    """
    main_trace = None
    if langfuse_client:
        main_trace = langfuse_client.trace(
            name="chat-rag-interaction-stream",
            user_id=user_id,
            session_id=session_id,
            input={"query": query},
            tags=["chat", "RAG", "streaming"],
        )

    response_sources = []

    try:
        if is_greeting(query):

            def greeting_stream():
                yield "Hello! How can I assist you today?"

            if main_trace:
                main_trace.update(
                    output={"response": "Hello! How can I assist you today?"}
                )
            return greeting_stream(), []

        # Use S3 vectors or local ChromaDB
        if config.USE_S3_VECTORS:
            try:
                retriever = create_s3_retriever(similarity_top_k=3)
                retrieved_nodes = retriever.retrieve(query)
            except Exception as e:
                err_msg = f"Error using S3 vectors: {str(e)}"
                if main_trace:
                    main_trace.update(output={"error": err_msg}, level="ERROR")

                def error_stream_s3():
                    yield err_msg

                return error_stream_s3(), []
        else:
            storage_context = get_storage_context()
            if storage_context is None:
                err_msg = "Error: Storage context not initialized. Cannot query index."
                if main_trace:
                    main_trace.update(output={"error": err_msg}, level="ERROR")

                def error_stream_storage():
                    yield err_msg

                return error_stream_storage(), []

            idx = load_index_from_storage(storage_context)
            retriever = idx.as_retriever(similarity_top_k=3)
            retrieved_nodes = retriever.retrieve(query)

        local_sources = []
        for n_ws in retrieved_nodes:
            node = n_ws.node
            source_text = (
                node.get_text()
                if hasattr(node, "get_text")
                else (node.text if hasattr(node, "text") else "")
            )

            # Check both file_path (ChromaDB) and source_file (S3)
            fp = node.metadata.get("file_path") or node.metadata.get("source_file")
            file_name = os.path.basename(fp) if fp else "Unknown Source"

            # Generate S3 document URL instead of local file path
            source_url = (
                f"/api/backend/files/s3-document?source_file={quote(file_name)}"
                if fp
                else None
            )

            local_sources.append(
                {
                    "file_name": file_name,
                    "file_path": fp,  # Keep for backwards compatibility
                    "source_url": source_url,  # New S3 URL for frontend
                    "page": node.metadata.get("page_number"),
                    "slide": node.metadata.get("slide_number"),
                    "chunk_text": source_text[:300] + "..."
                    if len(source_text) > 300
                    else source_text,
                }
            )
        nodes_for_prompt = retrieved_nodes
        # Format context with numbered sources for citation
        context_parts = []
        for idx, n_ws in enumerate(nodes_for_prompt, 1):
            source_text = (
                n_ws.node.get_text()
                if hasattr(n_ws.node, "get_text")
                else (n_ws.node.text if hasattr(n_ws.node, "text") else "")
            )
            source_name = local_sources[idx-1]["file_name"] if idx-1 < len(local_sources) else f"Source {idx}"
            context_parts.append(f"[Source {idx}: {source_name}]\n{source_text}")
        context_str = "\n\n".join(context_parts)
        response_sources = local_sources
        template_for_response = CHAT_QA_TEMPLATE
        used_web_search = False

        if _should_search_web(query, context_str, local_sources):
            web_results = search_web_results(query)
            if web_results:
                used_web_search = True
                template_for_response = WEB_SEARCH_TEMPLATE
                response_sources = []
                web_nodes = []
                for idx, result in enumerate(web_results):
                    snippet = result.get("content", "")
                    clean_snippet = (
                        snippet[:300] + "..." if len(snippet) > 300 else snippet
                    )
                    response_sources.append(
                        {
                            "file_name": result.get("title") or "Web result",
                            "file_path": None,
                            "external_url": result.get("url"),
                            "chunk_text": clean_snippet,
                            "source": result.get("source"),
                        }
                    )
                    text_payload = "\n".join(
                        filter(
                            None,
                            [
                                result.get("title"),
                                snippet,
                                f"Source: {result.get('url', '')}",
                            ],
                        )
                    )
                    node = TextNode(
                        id_=f"web-{idx}",
                        text=text_payload,
                        metadata={
                            "source": result.get("url"),
                            "title": result.get("title", "Web result"),
                        },
                    )
                    web_nodes.append(
                        NodeWithScore(
                            node=node,
                            score=float(result.get("score", 0.7) or 0.7),
                        )
                    )
                if web_nodes:
                    nodes_for_prompt = web_nodes
                    # Format web search context with numbered sources
                    web_context_parts = []
                    for idx, n in enumerate(nodes_for_prompt, 1):
                        source_title = n.node.metadata.get("title", f"Web Source {idx}")
                        web_context_parts.append(f"[Source {idx}: {source_title}]\n{n.node.get_text()}")
                    context_str = "\n\n".join(web_context_parts)
            else:
                used_web_search = False
                response_sources = local_sources
                template_for_response = CHAT_QA_TEMPLATE
                nodes_for_prompt = retrieved_nodes

        print("--------CONTEXT PASSED TO LLM--------")
        print(context_str)
        print("--------------------------------------")

        # Create response synthesizer with optional custom LLM
        synth_kwargs = {
            "response_mode": "compact",
            "streaming": True,
            "text_qa_template": template_for_response,
        }

        # If a specific model is requested, create a custom LLM instance
        if model_id:
            from backend.llm_provider import get_llm
            custom_llm = get_llm(model_id=model_id)
            synth_kwargs["llm"] = custom_llm
            logging.info(f"Using custom LLM model: {model_id}")

        synth = get_response_synthesizer(**synth_kwargs)
        streaming_response_obj = synth.synthesize(
            query=query,
            nodes=nodes_for_prompt,
        )
        response_text_generator = streaming_response_obj.response_gen

        if main_trace:
            main_trace.update(
                metadata={
                    "num_retrieved_sources": len(response_sources),
                    "retrieved_source_sample": [
                        s["file_name"] for s in response_sources[:2]
                    ],
                    "web_search_used": used_web_search,
                }
            )

        return response_text_generator, response_sources

    except Exception as e:
        logging.error(
            f"Error in generate_response_stream_and_sources: {e}", exc_info=True
        )
        error_message = str(e)
        if main_trace:
            main_trace.update(output={"error": error_message}, level="ERROR")

        def error_stream_exception():
            yield f"⚠️ Error processing your query: {error_message}"

        return error_stream_exception(), response_sources


def make_session_title(history: List) -> str:
    """Generate a title for chat session based on history"""
    if not history:
        return "New Chat"

    # Get the first user message for potential fallback
    first_user_message = None
    for entry in history:
        role, text = "", ""
        if isinstance(entry, (list, tuple)) and len(entry) >= 2:
            role, text = entry[0], entry[1]
        elif isinstance(entry, dict):
            role, text = entry.get("role", ""), entry.get("content", "")
        if role == "user" and text:
            first_user_message = text
            break

    # Build snippet from recent conversation
    snippet_lines = []
    for entry in history[-6:]:  # Last up to 6 messages
        role, text = "", ""
        if isinstance(entry, (list, tuple)) and len(entry) >= 2:
            role, text = entry[0], entry[1]
        elif isinstance(entry, dict):
            role, text = entry.get("role", ""), entry.get("content", "")

        if role and text:
            # Truncate very long messages for the prompt
            text_preview = text[:500] + "..." if len(text) > 500 else text
            snippet_lines.append(f"{role}: {text_preview}")

    if not snippet_lines:
        return (
            first_user_message.split()[:4].title() if first_user_message else "New Chat"
        )

    snippet = "\n".join(snippet_lines)

    try:
        prompt = (
            "Create a short, descriptive title (2-5 words) for this conversation:\n"
            f"{snippet}\n"
            "Just return the title, nothing else:"
        )
        title_response, _ = generate_response_with_sources(prompt)

        # Clean up the title
        title = title_response.strip()
        # Remove quotes, bullets, and extra whitespace
        title = re.sub(r'^["\'\-\*\s]+|["\'\-\*\s]+$', "", title)
        title = re.sub(r"\s+", " ", title).strip()

        # If title is too short or empty, use fallback
        if len(title) < 2:
            if first_user_message:
                return first_user_message.split()[:5].title()
            return "New Chat"

        # Limit to 5 words max
        words = title.split()[:5]
        title = " ".join(words).title()

        return (
            title
            if title
            else (
                first_user_message.split()[:5].title()
                if first_user_message
                else "New Chat"
            )
        )

    except Exception as e:
        logging.error(f"Error generating session title: {e}")
        # Fallback: use first few words of user's first message
        if first_user_message:
            words = first_user_message.split()
            # Take first 4 words, capitalizing each
            title_words = [w.capitalize() for w in words[:4]]
            return " ".join(title_words) if title_words else "New Chat"
        return "New Chat"


# --- File Conversion and Content Extraction ---
def convert_text_to_pdf(text: str, output_path: str) -> None:
    """Convert text to PDF file"""
    pdf = FPDF()
    pdf.add_page()

    # Try to use a Unicode font, fallback to Arial
    try:
        # You should replace this with a valid font path or use a bundled font
        font_path = "fonts/DejaVuSans.ttf"  # Relative path to avoid hardcoding
        if os.path.exists(font_path):
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", size=12)
        else:
            pdf.set_font("Arial", size=12)
    except Exception as e:
        logging.warning(f"Font loading failed, using Arial: {e}")
        pdf.set_font("Arial", size=12)

    pdf.set_auto_page_break(auto=True, margin=15)
    for line in text.split("\n"):
        try:
            pdf.multi_cell(0, 10, line)
        except UnicodeEncodeError:
            # Handle Unicode characters that can't be encoded
            safe_line = line.encode("latin-1", "ignore").decode("latin-1")
            pdf.multi_cell(0, 10, safe_line)

    pdf.output(output_path)


def convert_docx_to_pdf(docx_file_path: str, output_pdf_path: str) -> str:
    """Convert DOCX to PDF and return extracted text"""
    try:
        doc = DocxDoc(docx_file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        convert_text_to_pdf(text, output_pdf_path)
        return text
    except Exception as e:
        logging.error(f"Error converting DOCX to PDF: {e}")
        return f"Error processing DOCX file: {e}"


def convert_pptx_to_pdf(pptx_file_path: str, output_pdf_path: str) -> str:
    """Convert PPTX to PDF and return extracted text"""
    try:
        pres = Presentation(pptx_file_path)
        text = ""
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        convert_text_to_pdf(text, output_pdf_path)
        return text
    except Exception as e:
        logging.error(f"Error converting PPTX to PDF: {e}")
        return f"Error processing PPTX file: {e}"


def image_to_document(image_file_uploader_object) -> Document:
    """Convert image to document using OCR"""
    try:
        image = Image.open(image_file_uploader_object)
        text = pytesseract.image_to_string(image)
        return Document(text=text, metadata={"source": image_file_uploader_object.name})
    except pytesseract.TesseractNotFoundError:
        if st:
            st.error(
                "Tesseract is not installed or not in your PATH. OCR functionality will not work."
            )
        return Document(
            text="Error: Tesseract not found.",
            metadata={"source": image_file_uploader_object.name, "error": True},
        )
    except Exception as e:
        if st:
            st.error(f"Error during OCR: {e}")
        return Document(
            text=f"Error during OCR: {e}",
            metadata={"source": image_file_uploader_object.name, "error": True},
        )


def url_to_document(url: str) -> Document:
    """Fetch content from URL and convert to document"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()

        soup = BeautifulSoup(response.content, "html.parser")

        # Extract title
        page_title = url
        if soup.title and soup.title.string:
            page_title = soup.title.string.strip()

        # Remove unwanted elements
        for element in soup(["script", "style", "nav", "footer", "aside"]):
            element.decompose()

        text = soup.get_text(separator="\n", strip=True)

        if not text.strip():
            return Document(
                text=f"Warning: No main text content extracted from {url}. The page might be heavily JavaScript-reliant or empty.",
                metadata={"source": url, "title": page_title, "warning": True},
            )

        return Document(text=text, metadata={"source": url, "title": page_title})

    except requests.exceptions.RequestException as e:
        error_msg = f"Error fetching URL content from {url}: {e}"
        logging.error(error_msg)
        return Document(
            text=f"Error: {error_msg}",
            metadata={"source": url, "title": url, "error": True},
        )
    except Exception as e:
        error_msg = f"An unexpected error occurred while processing URL {url}: {e}"
        logging.error(error_msg)
        return Document(
            text=f"Error: {error_msg}",
            metadata={"source": url, "title": url, "error": True},
        )


def extract_video_id(url_or_id: str) -> Optional[str]:
    """Extract YouTube video ID from URL or validate existing ID"""
    if not url_or_id:
        return None

    # Regex to match various YouTube URL formats
    regex = r"(?:youtube\.com\/(?:[^\/\n\s]+\/\S+\/|(?:v|e(?:mbed)?)\/|\S*?[?&]v=)|youtu\.be\/|googleusercontent\.com\/youtube\.com\/\d+\/v\/)([a-zA-Z0-9_-]{11})"
    match = re.search(regex, url_or_id)

    if match:
        return match.group(1)

    # Check if it's already a valid video ID
    if re.fullmatch(r"[a-zA-Z0-9_-]{11}", url_or_id):
        return url_or_id

    return None


def youtube_to_document(video_id_or_url: str) -> Document:
    """Convert YouTube video to document using transcript"""
    video_id = extract_video_id(video_id_or_url)
    source_ref = video_id_or_url

    video_title = (
        f"YouTube Video: {video_id}"
        if video_id
        else f"Invalid YouTube Link: {video_id_or_url[:30]}..."
    )
    watch_url = (
        f"https://www.youtube.com/watch?v={video_id}" if video_id else source_ref
    )

    if not video_id:
        error_msg = f"Invalid YouTube URL or ID provided: '{video_id_or_url}'."
        logging.error(f"youtube_to_document: {error_msg}")
        return Document(
            text=f"Error: {error_msg}",
            metadata={"source": source_ref, "title": video_title, "error": True},
        )

    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript_obj = None

        # Try to get English transcript
        preferred_langs = ["en", "en-US", "en-GB"]
        try:
            transcript_obj = transcript_list.find_manually_created_transcript(
                preferred_langs
            )
        except NoTranscriptFound:
            try:
                transcript_obj = transcript_list.find_generated_transcript(
                    preferred_langs
                )
            except NoTranscriptFound:
                error_msg = f"No English transcript found for video ID '{video_id}'."
                logging.warning(f"youtube_to_document: {error_msg}")
                return Document(
                    text=f"Error: {error_msg}",
                    metadata={
                        "source": watch_url,
                        "title": video_title,
                        "video_id": video_id,
                        "error": True,
                    },
                )

        # Fetch transcript data
        fetched_transcript_data = transcript_obj.fetch()

        if not fetched_transcript_data:
            error_msg = f"Fetched transcript data is empty for video ID '{video_id}'."
            logging.warning(f"youtube_to_document: {error_msg}")
            return Document(
                text="Warning: Transcript was found but contained no segments.",
                metadata={
                    "source": watch_url,
                    "title": video_title,
                    "video_id": video_id,
                    "warning": True,
                },
            )

        # Extract text from transcript segments
        text_segments = []
        for segment in fetched_transcript_data:
            if isinstance(segment, dict) and "text" in segment:
                text_segments.append(segment["text"])
            elif hasattr(segment, "text"):
                text_segments.append(segment.text)

        full_transcript_text = " ".join(text_segments).strip()

        if not full_transcript_text:
            error_msg = f"Processed transcript for video ID '{video_id}' is empty."
            logging.warning(f"youtube_to_document: {error_msg}")
            return Document(
                text="Warning: Processed transcript was empty.",
                metadata={
                    "source": watch_url,
                    "title": video_title,
                    "video_id": video_id,
                    "warning": True,
                },
            )

        return Document(
            text=full_transcript_text,
            metadata={"source": watch_url, "title": video_title, "video_id": video_id},
        )

    except TranscriptsDisabled:
        error_msg = f"Transcripts are disabled for video ID '{video_id}'."
        logging.warning(f"youtube_to_document: {error_msg}")
        return Document(
            text=f"Error: {error_msg}",
            metadata={
                "source": watch_url,
                "title": video_title,
                "video_id": video_id,
                "error": True,
            },
        )
    except Exception as e:
        error_msg = f"Error processing transcript for video ID '{video_id}': {str(e)}"
        logging.error(f"youtube_to_document: {error_msg}", exc_info=True)
        return Document(
            text=f"Error: {error_msg}",
            metadata={
                "source": watch_url,
                "title": video_title,
                "video_id": video_id,
                "error": True,
            },
        )


def save_quiz_results(quiz_data: Dict[str, Any]) -> None:
    """Save quiz results to file"""
    folder_path = "quiz_results"
    os.makedirs(folder_path, exist_ok=True)

    now = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    filename = f"{folder_path}/quiz_{now}.json"

    try:
        with open(filename, "w", encoding="utf-8") as f:
            json.dump(quiz_data, f, indent=2, ensure_ascii=False)
        if st:
            st.success(f"Quiz results saved to `{filename}`")
    except Exception as e:
        if st:
            st.error(f"Failed to save quiz results: {e}")
        logging.error(f"Failed to save quiz results to {filename}: {e}")
