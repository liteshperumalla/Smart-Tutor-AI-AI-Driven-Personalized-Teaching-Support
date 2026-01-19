"""
Local rebuild: Process modules on local machine, upload to S3 with text + Bedrock
Runs outside Docker with more memory available
"""
import os
import sys
import json
import boto3
import pdfplumber
from pathlib import Path
from pptx import Presentation
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️  python-docx not available, .docx files will be skipped")

# Add project to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'backend'))

from backend.bedrock_embeddings import BedrockEmbeddings

print("="*80)
print("LOCAL S3 REBUILD WITH BEDROCK + TEXT")
print("="*80)

# Config
S3_BUCKET = "smart-tutor-vector-bucket"
S3_REGION = "us-east-1"
MODULES_DIR = "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules"
CHUNK_SIZE = 512

# Initialize
print("\n1. Initializing Bedrock...")
sys.stdout.flush()
bedrock = BedrockEmbeddings()
print("   ✅ Bedrock ready")
sys.stdout.flush()

print("\n   Initializing S3 client...")
sys.stdout.flush()
s3 = boto3.client('s3', region_name=S3_REGION)
print("   ✅ S3 client ready")
sys.stdout.flush()

# Text extraction functions
def extract_text_from_pdf(pdf_path):
    """Extract text from PDF"""
    texts = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
    except Exception as e:
        print(f"      ✗ PDF error: {e}")
    return "\n\n".join(texts)

def extract_text_from_pptx(pptx_path):
    """Extract text from PPTX/PPT"""
    texts = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                texts.append("\n".join(slide_text))
    except Exception as e:
        print(f"      ✗ PPTX/PPT error: {e}")
    return "\n\n".join(texts)

def extract_text_from_docx(docx_path):
    """Extract text from DOCX"""
    if not HAS_DOCX:
        return ""
    texts = []
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
    except Exception as e:
        print(f"      ✗ DOCX error: {e}")
    return "\n\n".join(texts)

def extract_text_from_ipynb(ipynb_path):
    """Extract text from Jupyter notebook"""
    texts = []
    try:
        with open(ipynb_path, 'r', encoding='utf-8') as f:
            notebook = json.load(f)
            for cell in notebook.get('cells', []):
                cell_type = cell.get('cell_type', '')
                source = cell.get('source', [])

                # Convert source to string
                if isinstance(source, list):
                    source = ''.join(source)

                if source.strip():
                    if cell_type == 'markdown':
                        texts.append(f"[Markdown]\n{source}")
                    elif cell_type == 'code':
                        texts.append(f"[Code]\n{source}")
    except Exception as e:
        print(f"      ✗ IPYNB error: {e}")
    return "\n\n".join(texts)

def simple_chunk(text, size=512):
    """Simple chunking by characters with overlap"""
    chunks = []
    overlap = size // 4  # 25% overlap
    for i in range(0, len(text), size - overlap):
        chunk = text[i:i+size]
        if len(chunk.strip()) > 50:  # Only keep meaningful chunks
            chunks.append(chunk)
    return chunks

# Process files
print(f"\n2. Processing files from {MODULES_DIR}...")
sys.stdout.flush()
uploaded = 0
failed_files = 0

if not os.path.exists(MODULES_DIR):
    print(f"\n❌ ERROR: Directory not found: {MODULES_DIR}")
    print("Please check the path and try again.")
    sys.exit(1)

print("   Counting files...")
sys.stdout.flush()
# Count total files first
total_files = sum(1 for root, dirs, files in os.walk(MODULES_DIR)
                  for file in files if file.endswith(('.pdf', '.pptx', '.ppt', '.docx', '.ipynb')))
print(f"   Found {total_files} files to process")
sys.stdout.flush()

processed_files = 0

print("   Starting file processing...")
sys.stdout.flush()

for root, dirs, files in os.walk(MODULES_DIR):
    for file in files:
        if file.endswith(('.pdf', '.pptx', '.ppt', '.docx', '.ipynb')):
            file_path = os.path.join(root, file)
            processed_files += 1
            print(f"\n   [{processed_files}/{total_files}] Processing: {file}")
            sys.stdout.flush()

            try:
                # Extract text based on file type
                if file.endswith('.pdf'):
                    text = extract_text_from_pdf(file_path)
                elif file.endswith(('.pptx', '.ppt')):
                    text = extract_text_from_pptx(file_path)
                elif file.endswith('.docx'):
                    text = extract_text_from_docx(file_path)
                elif file.endswith('.ipynb'):
                    text = extract_text_from_ipynb(file_path)

                if not text or len(text.strip()) < 100:
                    print(f"      → Skipped (no text)")
                    continue

                # Chunk
                chunks = simple_chunk(text, CHUNK_SIZE)
                print(f"      → {len(chunks)} chunks created")

                # Upload each chunk
                for i, chunk_text in enumerate(chunks):
                    try:
                        # Generate embedding
                        embedding = bedrock._encode_single(chunk_text)

                        # Create chunk
                        chunk_id = f"{Path(file).stem}_chunk_{i:04d}"
                        chunk_data = {
                            'chunk_id': chunk_id,
                            'text': chunk_text,  # ← KEY: Include text!
                            'embedding': embedding,
                            'source_file': file,
                            'chunk_index': i
                        }

                        # Upload to S3
                        s3.put_object(
                            Bucket=S3_BUCKET,
                            Key=f"chunks/{chunk_id}.json",
                            Body=json.dumps(chunk_data),
                            ContentType='application/json'
                        )

                        uploaded += 1

                        if uploaded % 50 == 0:
                            print(f"      → Progress: {uploaded} chunks uploaded...")

                    except Exception as e:
                        print(f"      ✗ Chunk {i} error: {e}")

            except Exception as e:
                print(f"      ✗ File error: {e}")
                failed_files += 1

print(f"\n{'='*80}")
print(f"✅ COMPLETE!")
print(f"{'='*80}")
print(f"Uploaded: {uploaded} chunks")
print(f"Failed files: {failed_files}")
print(f"Location: s3://{S3_BUCKET}/chunks/")
print(f"\nNext: Restart backend to use the new chunks")
