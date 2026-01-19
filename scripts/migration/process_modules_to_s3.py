#!/usr/bin/env python3
"""
Process course modules and upload to S3 with Bedrock embeddings
Follows exact chunking strategy from Data_parsing.py and Data_loading.py
"""

import os
import sys
import re
import json
import time
import hashlib
import unicodedata
import boto3
import pdfplumber
from pathlib import Path
from pptx import Presentation
from typing import List, Dict

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from backend.config import config
from backend.bedrock_embeddings import BedrockEmbeddings
from backend.logger import get_logger
from llama_index.core import SimpleDirectoryReader, Settings
from llama_index.core.schema import Document
from llama_index.core.node_parser import SentenceSplitter, SemanticSplitterNodeParser
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core.readers.base import BaseReader

logger = get_logger(__name__)

# Configuration
BUCKET_NAME = "smart-ai-tutor-docs"
REGION = config.AWS_REGION
MODULES_DIR = "./Modules"

# Chunking configuration (matching Data_parsing.py line 321-324)
CHUNK_SIZE = 512  # characters
CHUNK_OVERLAP = 102  # 20% overlap

# Initialize Bedrock embeddings
bedrock_embeddings = BedrockEmbeddings(
    model_id=config.BEDROCK_EMBEDDING_MODEL_ID,
    region=REGION
)

# Initialize HuggingFace embedding for semantic splitter (matching Data_parsing.py line 452)
try:
    model_name = "BAAI/bge-small-en-v1.5"
    Settings.embed_model = HuggingFaceEmbedding(model_name=model_name)
    print(f"✅ Semantic splitter model loaded: {model_name}")
except Exception as e:
    print(f"⚠️ Falling back to all-MiniLM-L6-v2: {e}")
    model_name = "sentence-transformers/all-MiniLM-L6-v2"
    Settings.embed_model = HuggingFaceEmbedding(model_name=model_name)

# Initialize S3 client
s3_client = boto3.client('s3', region_name=REGION)

# ------------------------------
# TEXT PREPROCESSING (from Data_parsing.py lines 471-508)
# ------------------------------
def clean_text(text):
    """Clean text while preserving important information"""
    text = unicodedata.normalize("NFKD", text)

    # Preserve email addresses
    email_pattern = r'([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})'
    emails = re.findall(email_pattern, text)
    for i, email in enumerate(emails):
        text = text.replace(email, f'EMAIL_PLACEHOLDER_{i}')

    # Preserve URLs
    url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
    urls = re.findall(url_pattern, text)
    for i, url in enumerate(urls):
        text = text.replace(url, f'URL_PLACEHOLDER_{i}')

    # Remove table of contents artifacts
    text = re.sub(r'\.{5,}', ' ', text)

    # Remove excessive spaces
    text = re.sub(r'\s+', ' ', text).strip()

    # Restore emails and URLs
    for i, email in enumerate(emails):
        text = text.replace(f'EMAIL_PLACEHOLDER_{i}', email)
    for i, url in enumerate(urls):
        text = text.replace(f'URL_PLACEHOLDER_{i}', url)

    return text

def preprocess_text(file_path, text):
    """Preprocess text based on file type"""
    code_extensions = {".py", ".java", ".cpp", ".js", ".c", ".cs", ".html", ".css", ".php", ".rb", ".ipynb"}
    text_extensions = {".pdf", ".docx", ".pptx", ".txt"}

    ext = os.path.splitext(file_path)[-1].lower()

    if ext in text_extensions or ext in code_extensions:
        text = clean_text(text)
    return text

# ------------------------------
# CUSTOM READERS (from Data_parsing.py lines 35-183)
# ------------------------------
class PPTXTextOnlyReader(BaseReader):
    """Custom PPTX reader that extracts text slide-by-slide"""
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        prs = Presentation(file_path)
        docs: list[Document] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw = shape.text.strip()
                    if raw:
                        text_runs.append(preprocess_text(file_path, raw))
            slide_text = "\n".join(text_runs)
            if slide_text:
                docs.append(Document(
                    text=slide_text,
                    metadata={
                        "file_path": file_path,
                        "file_name": os.path.basename(file_path),
                        "slide_number": slide_idx,
                        "file_type": "pptx",
                        "folder_name": os.path.basename(os.path.dirname(file_path))
                    }
                ))
        return docs

class PDFTextOnlyReader(BaseReader):
    """Custom PDF reader that extracts text page-by-page"""
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        docs: list[Document] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                raw = page.extract_text()
                if raw:
                    cleaned = preprocess_text(file_path, raw)
                    docs.append(Document(
                        text=cleaned,
                        metadata={
                            "file_path": file_path,
                            "file_name": os.path.basename(file_path),
                            "page_number": page.page_number,
                            "file_type": "pdf",
                            "folder_name": os.path.basename(os.path.dirname(file_path))
                        }
                    ))
        return docs

class NotebookReader(BaseReader):
    """Custom notebook reader"""
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                notebook = json.load(f)

            text_parts = []
            for cell in notebook.get('cells', []):
                cell_type = cell.get('cell_type', '')
                if cell_type in ['markdown', 'code']:
                    source = cell.get('source', [])
                    if isinstance(source, list):
                        text_parts.append(''.join(source))
                    else:
                        text_parts.append(source)

            full_text = '\n\n'.join(text_parts)

            return [Document(
                text=preprocess_text(file_path, full_text),
                metadata={
                    'file_path': file_path,
                    'file_name': os.path.basename(file_path),
                    'file_type': 'notebook',
                    'folder_name': os.path.basename(os.path.dirname(file_path))
                }
            )]
        except Exception as e:
            logger.error(f"Error reading notebook {file_path}: {e}")
            return []

# ------------------------------
# CONTEXTUAL ENRICHMENT (from Data_parsing.py lines 235-283)
# ------------------------------
def enrich_chunk_with_context(text, metadata):
    """Add contextual information to chunk"""
    context_parts = []

    if 'file_name' in metadata:
        context_parts.append(f"Document: {metadata['file_name']}")

    if 'folder_name' in metadata:
        context_parts.append(f"Section: {metadata['folder_name']}")

    location_info = []
    if 'page_number' in metadata:
        location_info.append(f"Page {metadata['page_number']}")
    elif 'slide_number' in metadata:
        location_info.append(f"Slide {metadata['slide_number']}")

    if location_info:
        context_parts.append(f"Location: {', '.join(location_info)}")

    if 'file_type' in metadata:
        file_type = metadata['file_type']
        if file_type == 'notebook':
            context_parts.append("Source: Jupyter Notebook")
        elif file_type == 'pdf':
            context_parts.append("Source: PDF Document")
        elif file_type == 'pptx':
            context_parts.append("Source: PowerPoint Presentation")

    if context_parts:
        context_header = " | ".join(context_parts)
        enriched_text = f"[CONTEXT: {context_header}]\n\n{text}"
        return enriched_text

    return text

# ------------------------------
# NOTEBOOK-AWARE PARSER (from Data_parsing.py lines 288-443)
# ------------------------------
class NotebookAwareParser:
    """Content-type aware parser matching Data_parsing.py"""

    def __init__(self, embed_model):
        self.embed_model = embed_model

        # Semantic splitter for documents
        self.semantic_splitter = SemanticSplitterNodeParser(
            buffer_size=1,
            breakpoint_percentile_threshold=95,
            embed_model=embed_model
        )

        # Sentence splitter (matching Data_parsing.py line 321-324)
        self.sentence_splitter = SentenceSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP
        )

    def parse_documents(self, documents):
        """Parse documents using appropriate parser"""
        all_nodes = []

        for doc in documents:
            # Apply contextual enrichment
            enriched_text = enrich_chunk_with_context(doc.text, doc.metadata)
            doc = Document(text=enriched_text, metadata=doc.metadata.copy())

            file_type = doc.metadata.get('file_type', 'unknown')

            # Parse based on file type
            if file_type == 'notebook':
                nodes = self._parse_notebook(doc)
            else:
                nodes = self._parse_regular_document(doc)

            all_nodes.extend(nodes)

        return all_nodes

    def _parse_notebook(self, doc):
        """Parse notebook with semantic splitter"""
        try:
            return self.semantic_splitter.get_nodes_from_documents([doc])
        except Exception as e:
            logger.warning(f"Semantic parsing failed, using sentence splitter: {e}")
            return self.sentence_splitter.get_nodes_from_documents([doc])

    def _parse_regular_document(self, doc):
        """Parse regular documents with semantic splitter"""
        try:
            return self.semantic_splitter.get_nodes_from_documents([doc])
        except Exception as e:
            logger.warning(f"Semantic parsing failed, using sentence splitter: {e}")
            return self.sentence_splitter.get_nodes_from_documents([doc])

# ------------------------------
# MAIN PROCESSING PIPELINE
# ------------------------------
class ModulesToS3Processor:
    """Process modules and upload to S3 with Bedrock embeddings"""

    def __init__(self):
        self.s3 = s3_client
        self.embeddings = bedrock_embeddings
        self.parser = NotebookAwareParser(Settings.embed_model)
        self.uploaded_count = 0
        self.total_cost = 0.0

    def process_modules(self, modules_dir: str):
        """Process all modules from directory"""
        print("=" * 70)
        print("PROCESSING MODULES TO S3")
        print("=" * 70)
        print(f"Modules directory: {modules_dir}")
        print(f"S3 bucket: {BUCKET_NAME}")
        print(f"Chunk size: {CHUNK_SIZE} characters")
        print(f"Chunk overlap: {CHUNK_OVERLAP} characters")
        print(f"Embedding model: {config.BEDROCK_EMBEDDING_MODEL_ID}")
        print("=" * 70)
        print()

        # Load documents from Modules directory
        print(f"📁 Loading documents from {modules_dir}...")

        try:
            reader = SimpleDirectoryReader(
                input_dir=modules_dir,
                required_exts=['.pptx', '.ipynb', '.docx', '.pdf', '.txt', '.md'],
                file_extractor={
                    ".pptx": PPTXTextOnlyReader(),
                    ".pdf": PDFTextOnlyReader(),
                    ".ipynb": NotebookReader()
                },
                recursive=True
            )
            docs = reader.load_data()

            if not docs:
                print("❌ No documents found!")
                return

            print(f"✅ Loaded {len(docs)} documents")
            print()

        except Exception as e:
            print(f"❌ Error loading documents: {e}")
            return

        # Parse documents with notebook-aware parser
        print("🔄 Parsing documents with content-aware chunking...")
        nodes = self.parser.parse_documents(docs)
        print(f"✅ Created {len(nodes)} chunks")
        print()

        # Generate embeddings and upload
        print("🚀 Generating embeddings and uploading to S3...")
        self._process_and_upload_chunks(nodes)

        print()
        print("=" * 70)
        print("✅ PROCESSING COMPLETE!")
        print(f"   Uploaded: {self.uploaded_count} chunks")
        print(f"   Estimated cost: ${self.total_cost:.4f}")
        print("=" * 70)

    def _process_and_upload_chunks(self, nodes):
        """Generate embeddings and upload chunks to S3"""
        batch_size = 10  # Process in batches

        for i in range(0, len(nodes), batch_size):
            batch = nodes[i:i + batch_size]
            print(f"  Processing batch {i//batch_size + 1}/{(len(nodes) + batch_size - 1)//batch_size}...")

            for node in batch:
                try:
                    self._process_single_chunk(node)
                except Exception as e:
                    logger.error(f"Error processing chunk: {e}")

            # Small delay between batches
            if i + batch_size < len(nodes):
                time.sleep(0.5)

    def _process_single_chunk(self, node):
        """Process and upload a single chunk"""
        # Generate embedding
        text = node.get_content() if hasattr(node, 'get_content') else node.text

        try:
            embedding = self.embeddings.encode([text])[0]
        except Exception as e:
            logger.error(f"Error generating embedding: {e}")
            return

        # Create chunk ID
        file_path = node.metadata.get('file_path', '')
        relative_path = os.path.relpath(file_path, MODULES_DIR) if file_path else 'unknown'

        file_hash = hashlib.sha256(relative_path.encode('utf-8')).hexdigest()[:8]
        chunk_index = node.metadata.get('chunk_index', self.uploaded_count)
        chunk_id = f"{Path(relative_path).stem}_{file_hash}_chunk_{chunk_index:04d}"

        # S3 keys
        chunk_key = f"chunks/{relative_path}/chunk_{chunk_index:04d}.txt"
        vector_key = f"chunks/{relative_path}/chunk_{chunk_index:04d}.vector.json"

        # Upload chunk text
        try:
            self.s3.put_object(
                Bucket=BUCKET_NAME,
                Key=chunk_key,
                Body=text.encode('utf-8'),
                ContentType='text/plain',
                Metadata={
                    'chunk-id': chunk_id,
                    'source-file': relative_path,
                    'chunk-index': str(chunk_index)
                }
            )

            # Upload vector
            vector_data = {
                'chunk_id': chunk_id,
                'embedding': embedding,
                'dimension': len(embedding),
                'source_file': relative_path,
                'chunk_index': chunk_index
            }

            self.s3.put_object(
                Bucket=BUCKET_NAME,
                Key=vector_key,
                Body=json.dumps(vector_data).encode('utf-8'),
                ContentType='application/json'
            )

            self.uploaded_count += 1

            # Estimate cost (Titan: $0.0001 per 1K tokens)
            tokens = len(text.split())
            self.total_cost += (tokens / 1000) * 0.0001

        except Exception as e:
            logger.error(f"Failed to upload chunk {chunk_id}: {e}")

# ------------------------------
# MAIN ENTRY POINT
# ------------------------------
def main():
    processor = ModulesToS3Processor()
    processor.process_modules(MODULES_DIR)

if __name__ == "__main__":
    main()
