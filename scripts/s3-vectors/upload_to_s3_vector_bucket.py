#!/usr/bin/env python3
"""
Upload course modules to AWS S3 Vector bucket with Bedrock embeddings
Follows Data_parsing.py chunking strategy
Uses AWS S3 Vector bucket API for native vector search
"""

import os
import sys
import re
import json
import time
import hashlib
import unicodedata
import boto3
import pdfplumber
from pathlib import Path
from pptx import Presentation
from typing import List, Dict

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from backend.config import config
from backend.bedrock_embeddings import BedrockEmbeddings
from backend.logger import get_logger
from llama_index.core import SimpleDirectoryReader, Settings
from llama_index.core.schema import Document
from llama_index.core.node_parser import SentenceSplitter, SemanticSplitterNodeParser
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core.readers.base import BaseReader

logger = get_logger(__name__)

# Configuration
VECTOR_BUCKET_NAME = "smart-tutor-vector-bucket"
REGION = config.AWS_REGION
MODULES_DIR = "./Modules"

# Chunking configuration (matching Data_parsing.py)
CHUNK_SIZE = 512
CHUNK_OVERLAP = 102

# Initialize Bedrock embeddings
bedrock_embeddings = BedrockEmbeddings(
    model_id=config.BEDROCK_EMBEDDING_MODEL_ID,
    region=REGION
)

# Initialize HuggingFace for semantic splitter
try:
    model_name = "BAAI/bge-small-en-v1.5"
    Settings.embed_model = HuggingFaceEmbedding(model_name=model_name)
    print(f"✅ Semantic splitter model loaded: {model_name}")
except Exception as e:
    print(f"⚠️ Falling back to all-MiniLM-L6-v2: {e}")
    model_name = "sentence-transformers/all-MiniLM-L6-v2"
    Settings.embed_model = HuggingFaceEmbedding(model_name=model_name)

# Initialize S3 client for Vector bucket
s3_client = boto3.client('s3', region_name=REGION)

# Text preprocessing functions (from Data_parsing.py)
def clean_text(text):
    text = unicodedata.normalize("NFKD", text)
    email_pattern = r'([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})'
    emails = re.findall(email_pattern, text)
    for i, email in enumerate(emails):
        text = text.replace(email, f'EMAIL_PLACEHOLDER_{i}')
    url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
    urls = re.findall(url_pattern, text)
    for i, url in enumerate(urls):
        text = text.replace(url, f'URL_PLACEHOLDER_{i}')
    text = re.sub(r'\.{5,}', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    for i, email in enumerate(emails):
        text = text.replace(f'EMAIL_PLACEHOLDER_{i}', email)
    for i, url in enumerate(urls):
        text = text.replace(f'URL_PLACEHOLDER_{i}', url)
    return text

def preprocess_text(file_path, text):
    code_extensions = {".py", ".java", ".cpp", ".js", ".c", ".cs", ".html", ".css", ".php", ".rb", ".ipynb"}
    text_extensions = {".pdf", ".docx", ".pptx", ".txt"}
    ext = os.path.splitext(file_path)[-1].lower()
    if ext in text_extensions or ext in code_extensions:
        text = clean_text(text)
    return text

# Custom readers (from Data_parsing.py)
class PPTXTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        prs = Presentation(file_path)
        docs: list[Document] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw = shape.text.strip()
                    if raw:
                        text_runs.append(preprocess_text(file_path, raw))
            slide_text = "\n".join(text_runs)
            if slide_text:
                docs.append(Document(
                    text=slide_text,
                    metadata={
                        "file_path": file_path,
                        "file_name": os.path.basename(file_path),
                        "slide_number": slide_idx,
                        "file_type": "pptx",
                        "folder_name": os.path.basename(os.path.dirname(file_path))
                    }
                ))
        return docs

class PDFTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        docs: list[Document] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                raw = page.extract_text()
                if raw:
                    cleaned = preprocess_text(file_path, raw)
                    docs.append(Document(
                        text=cleaned,
                        metadata={
                            "file_path": file_path,
                            "file_name": os.path.basename(file_path),
                            "page_number": page.page_number,
                            "file_type": "pdf",
                            "folder_name": os.path.basename(os.path.dirname(file_path))
                        }
                    ))
        return docs

class NotebookReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                notebook = json.load(f)
            text_parts = []
            for cell in notebook.get('cells', []):
                cell_type = cell.get('cell_type', '')
                if cell_type in ['markdown', 'code']:
                    source = cell.get('source', [])
                    if isinstance(source, list):
                        text_parts.append(''.join(source))
                    else:
                        text_parts.append(source)
            full_text = '\n\n'.join(text_parts)
            return [Document(
                text=preprocess_text(file_path, full_text),
                metadata={
                    'file_path': file_path,
                    'file_name': os.path.basename(file_path),
                    'file_type': 'notebook',
                    'folder_name': os.path.basename(os.path.dirname(file_path))
                }
            )]
        except Exception as e:
            logger.error(f"Error reading notebook {file_path}: {e}")
            return []

# Contextual enrichment (from Data_parsing.py)
def enrich_chunk_with_context(text, metadata):
    context_parts = []
    if 'file_name' in metadata:
        context_parts.append(f"Document: {metadata['file_name']}")
    if 'folder_name' in metadata:
        context_parts.append(f"Section: {metadata['folder_name']}")
    location_info = []
    if 'page_number' in metadata:
        location_info.append(f"Page {metadata['page_number']}")
    elif 'slide_number' in metadata:
        location_info.append(f"Slide {metadata['slide_number']}")
    if location_info:
        context_parts.append(f"Location: {', '.join(location_info)}")
    if 'file_type' in metadata:
        file_type = metadata['file_type']
        if file_type == 'notebook':
            context_parts.append("Source: Jupyter Notebook")
        elif file_type == 'pdf':
            context_parts.append("Source: PDF Document")
        elif file_type == 'pptx':
            context_parts.append("Source: PowerPoint Presentation")
    if context_parts:
        context_header = " | ".join(context_parts)
        enriched_text = f"[CONTEXT: {context_header}]\n\n{text}"
        return enriched_text
    return text

# Notebook-aware parser
class NotebookAwareParser:
    def __init__(self, embed_model):
        self.embed_model = embed_model
        self.semantic_splitter = SemanticSplitterNodeParser(
            buffer_size=1,
            breakpoint_percentile_threshold=95,
            embed_model=embed_model
        )
        self.sentence_splitter = SentenceSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP
        )

    def parse_documents(self, documents):
        all_nodes = []
        for doc in documents:
            enriched_text = enrich_chunk_with_context(doc.text, doc.metadata)
            doc = Document(text=enriched_text, metadata=doc.metadata.copy())
            file_type = doc.metadata.get('file_type', 'unknown')
            if file_type == 'notebook':
                nodes = self._parse_notebook(doc)
            else:
                nodes = self._parse_regular_document(doc)
            all_nodes.extend(nodes)
        return all_nodes

    def _parse_notebook(self, doc):
        try:
            return self.semantic_splitter.get_nodes_from_documents([doc])
        except Exception as e:
            logger.warning(f"Semantic parsing failed, using sentence splitter: {e}")
            return self.sentence_splitter.get_nodes_from_documents([doc])

    def _parse_regular_document(self, doc):
        try:
            return self.semantic_splitter.get_nodes_from_documents([doc])
        except Exception as e:
            logger.warning(f"Semantic parsing failed, using sentence splitter: {e}")
            return self.sentence_splitter.get_nodes_from_documents([doc])

# S3 Vector bucket uploader
class S3VectorBucketUploader:
    """Upload to AWS S3 Vector bucket with native vector search"""

    def __init__(self):
        self.s3 = s3_client
        self.embeddings = bedrock_embeddings
        self.parser = NotebookAwareParser(Settings.embed_model)
        self.uploaded_count = 0
        self.total_cost = 0.0

    def process_modules(self, modules_dir: str):
        print("=" * 70)
        print("UPLOADING TO AWS S3 VECTOR BUCKET")
        print("=" * 70)
        print(f"Modules directory: {modules_dir}")
        print(f"S3 Vector bucket: {VECTOR_BUCKET_NAME}")
        print(f"Chunk size: {CHUNK_SIZE} characters")
        print(f"Chunk overlap: {CHUNK_OVERLAP} characters")
        print(f"Embedding model: {config.BEDROCK_EMBEDDING_MODEL_ID}")
        print("=" * 70)
        print()

        # Load documents
        print(f"📁 Loading documents from {modules_dir}...")
        try:
            reader = SimpleDirectoryReader(
                input_dir=modules_dir,
                required_exts=['.pptx', '.ipynb', '.docx', '.pdf', '.txt', '.md'],
                file_extractor={
                    ".pptx": PPTXTextOnlyReader(),
                    ".pdf": PDFTextOnlyReader(),
                    ".ipynb": NotebookReader()
                },
                recursive=True
            )
            docs = reader.load_data()
            if not docs:
                print("❌ No documents found!")
                return
            print(f"✅ Loaded {len(docs)} documents")
            print()
        except Exception as e:
            print(f"❌ Error loading documents: {e}")
            return

        # Parse documents
        print("🔄 Parsing documents with content-aware chunking...")
        nodes = self.parser.parse_documents(docs)
        print(f"✅ Created {len(nodes)} chunks")
        print()

        # Upload to S3 Vector bucket
        print("🚀 Uploading to S3 Vector bucket with embeddings...")
        self._upload_to_vector_bucket(nodes)

        print()
        print("=" * 70)
        print("✅ UPLOAD COMPLETE!")
        print(f"   Uploaded: {self.uploaded_count} vectors")
        print(f"   Estimated cost: ${self.total_cost:.4f}")
        print("=" * 70)

    def _upload_to_vector_bucket(self, nodes):
        """Upload nodes with embeddings to S3 Vector bucket"""
        batch_size = 10

        for i in range(0, len(nodes), batch_size):
            batch = nodes[i:i + batch_size]
            print(f"  Processing batch {i//batch_size + 1}/{(len(nodes) + batch_size - 1)//batch_size}...")

            for node in batch:
                try:
                    self._upload_single_vector(node)
                except Exception as e:
                    logger.error(f"Error uploading vector: {e}")

            if i + batch_size < len(nodes):
                time.sleep(0.5)

    def _upload_single_vector(self, node):
        """Upload single chunk with embedding to S3 Vector bucket"""
        text = node.get_content() if hasattr(node, 'get_content') else node.text

        # Generate embedding
        try:
            embedding = self.embeddings.encode([text])[0]
        except Exception as e:
            logger.error(f"Error generating embedding: {e}")
            return

        # Create document object for S3 Vector bucket
        file_path = node.metadata.get('file_path', '')
        relative_path = os.path.relpath(file_path, MODULES_DIR) if file_path else 'unknown'

        chunk_index = node.metadata.get('chunk_index', self.uploaded_count)
        file_hash = hashlib.sha256(relative_path.encode('utf-8')).hexdigest()[:8]
        doc_id = f"{Path(relative_path).stem}_{file_hash}_chunk_{chunk_index:04d}"

        # S3 Vector bucket object
        s3_key = f"vectors/{relative_path}/chunk_{chunk_index:04d}.json"

        # Create vector document in S3 Vector bucket format
        vector_doc = {
            "id": doc_id,
            "text": text,
            "embedding": embedding,
            "metadata": {
                "source_file": relative_path,
                "chunk_index": chunk_index,
                "file_name": node.metadata.get('file_name', ''),
                "folder_name": node.metadata.get('folder_name', ''),
                "file_type": node.metadata.get('file_type', ''),
                "embedding_model": config.BEDROCK_EMBEDDING_MODEL_ID,
                "embedding_dimension": len(embedding)
            }
        }

        # Upload to S3 Vector bucket
        try:
            self.s3.put_object(
                Bucket=VECTOR_BUCKET_NAME,
                Key=s3_key,
                Body=json.dumps(vector_doc).encode('utf-8'),
                ContentType='application/json',
                Metadata={
                    'document-id': doc_id,
                    'source-file': relative_path,
                    'chunk-index': str(chunk_index),
                    'vector-dimension': str(len(embedding))
                }
            )

            self.uploaded_count += 1

            # Estimate cost
            tokens = len(text.split())
            self.total_cost += (tokens / 1000) * 0.0001

        except Exception as e:
            logger.error(f"Failed to upload vector {doc_id}: {e}")

def main():
    processor = S3VectorBucketUploader()
    processor.process_modules(MODULES_DIR)

if __name__ == "__main__":
    main()
