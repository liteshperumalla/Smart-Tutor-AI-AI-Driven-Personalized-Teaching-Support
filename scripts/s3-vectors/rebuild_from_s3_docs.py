#!/usr/bin/env python3
"""
Rebuild S3 Vector Index from documents in S3 bucket
Downloads all files from s3://smart-ai-tutor-docs/modules/
Processes them, creates chunks with embeddings, and uploads to S3
"""

import os
import sys
import json
import boto3
import pdfplumber
import pickle
import numpy as np
from pathlib import Path
from pptx import Presentation
import hashlib
import re

try:
    from docx import Document

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️  python-docx not available, .docx files will be skipped")

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "backend"))

from backend.bedrock_embeddings import BedrockEmbeddings

print("=" * 80)
print("S3 VECTOR INDEX REBUILD")
print("=" * 80)

# Configuration
S3_BUCKET = "smart-ai-tutor-docs"
S3_REGION = "us-east-1"
DOCS_PREFIX = "modules/"
CHUNKS_PREFIX = "chunks/"
INDEX_KEY = "vector_index/s3_vector_index.pkl"
CHUNK_SIZE = 512
CHUNK_OVERLAP = 128


def extract_text_from_pdf(pdf_content):
    """Extract text from PDF bytes"""
    texts = []
    try:
        import io

        with pdfplumber.open(io.BytesIO(pdf_content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
    except Exception as e:
        print(f"      ✗ PDF error: {e}")
    return "\n\n".join(texts)


def extract_text_from_pptx(pptx_content):
    """Extract text from PPTX/PPT bytes"""
    texts = []
    try:
        import io

        prs = Presentation(io.BytesIO(pptx_content))
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                texts.append("\n".join(slide_text))
    except Exception as e:
        print(f"      ✗ PPTX/PPT error: {e}")
    return "\n\n".join(texts)


def extract_text_from_docx(docx_content):
    """Extract text from DOCX bytes"""
    if not HAS_DOCX:
        return ""
    texts = []
    try:
        import io

        doc = Document(io.BytesIO(docx_content))
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
    except Exception as e:
        print(f"      ✗ DOCX error: {e}")
    return "\n\n".join(texts)


def extract_text_from_ipynb(ipynb_content):
    """Extract text from Jupyter notebook bytes"""
    texts = []
    try:
        notebook = json.loads(ipynb_content.decode("utf-8"))
        for cell in notebook.get("cells", []):
            cell_type = cell.get("cell_type", "")
            source = cell.get("source", [])
            if isinstance(source, list):
                source = "".join(source)
            if source.strip():
                if cell_type == "markdown":
                    texts.append(f"[Markdown]\n{source}")
                elif cell_type == "code":
                    texts.append(f"[Code]\n{source}")
    except Exception as e:
        print(f"      ✗ IPYNB error: {e}")
    return "\n\n".join(texts)


def simple_chunk(text, size=512, overlap=128):
    """Simple chunking by characters with overlap"""
    chunks = []
    for i in range(0, len(text), size - overlap):
        chunk = text[i : i + size]
        if len(chunk.strip()) > 50:
            chunks.append(chunk)
    return chunks


def normalize_path(path):
    """Normalize path to handle case sensitivity issues"""
    return path


def main():
    # Initialize
    print("\n1. Initializing Bedrock embeddings...")
    sys.stdout.flush()
    bedrock = BedrockEmbeddings()
    print("   ✅ Bedrock ready")
    sys.stdout.flush()

    print("\n2. Initializing S3 client...")
    sys.stdout.flush()
    s3 = boto3.client("s3", region_name=S3_REGION)
    print("   ✅ S3 client ready")
    sys.stdout.flush()

    # List all files in S3
    print(f"\n3. Scanning S3 bucket: s3://{S3_BUCKET}/{DOCS_PREFIX}")
    sys.stdout.flush()

    all_files = []
    paginator = s3.get_paginator("list_objects_v2")

    for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=DOCS_PREFIX):
        if "Contents" not in page:
            continue
        for obj in page["Contents"]:
            key = obj["Key"]
            if key.endswith((".pdf", ".pptx", ".ppt", ".docx", ".ipynb")):
                all_files.append(key)

    print(f"   ✅ Found {len(all_files)} files to process")
    sys.stdout.flush()

    # Process files
    print(f"\n4. Processing files and creating chunks...")
    sys.stdout.flush()

    all_vectors = []
    all_metadata = []
    total_chunks = 0
    processed_files = 0
    failed_files = 0

    for file_key in all_files:
        processed_files += 1
        filename = file_key.split("/")[-1]
        print(f"\n   [{processed_files}/{len(all_files)}] {filename}")
        sys.stdout.flush()

        try:
            # Download file
            response = s3.get_object(Bucket=S3_BUCKET, Key=file_key)
            content = response["Body"].read()

            # Extract text
            if filename.endswith(".pdf"):
                text = extract_text_from_pdf(content)
            elif filename.endswith((".pptx", ".ppt")):
                text = extract_text_from_pptx(content)
            elif filename.endswith(".docx"):
                text = extract_text_from_docx(content)
            elif filename.endswith(".ipynb"):
                text = extract_text_from_ipynb(content)
            else:
                text = ""

            if not text or len(text.strip()) < 100:
                print(f"      → Skipped (no text extracted)")
                continue

            # Create chunks
            chunks = simple_chunk(text, CHUNK_SIZE, CHUNK_OVERLAP)
            print(f"      → {len(chunks)} chunks")

            # Use relative path from modules/ as source_file
            source_file = file_key.replace(DOCS_PREFIX, "", 1)

            # Process each chunk
            for i, chunk_text in enumerate(chunks):
                try:
                    # Generate embedding
                    embedding = bedrock._encode_single(chunk_text)

                    # Create unique chunk ID
                    file_stem = Path(filename).stem
                    # Sanitize filename for chunk_id
                    safe_stem = re.sub(r"[^a-zA-Z0-9_-]", "_", file_stem)[:50]
                    chunk_id = f"{safe_stem}_chunk_{total_chunks:05d}"

                    # Create chunk data with correct source_file path
                    chunk_data = {
                        "chunk_id": chunk_id,
                        "text": chunk_text,
                        "embedding": embedding,
                        "source_file": source_file,
                        "chunk_index": i,
                        "file_path": source_file,
                    }

                    # Upload chunk to S3
                    chunk_key = f"{CHUNKS_PREFIX}{source_file}/{chunk_id}.json"
                    s3.put_object(
                        Bucket=S3_BUCKET,
                        Key=chunk_key,
                        Body=json.dumps(chunk_data),
                        ContentType="application/json",
                    )

                    # Store for index
                    all_vectors.append(embedding)
                    all_metadata.append(
                        {
                            "chunk_id": chunk_id,
                            "source_file": source_file,
                            "file_path": source_file,
                            "chunk_index": i,
                            "s3_key": chunk_key.replace(".json", ".txt"),
                        }
                    )

                    total_chunks += 1

                    if total_chunks % 100 == 0:
                        print(f"      → Progress: {total_chunks} chunks...")

                except Exception as e:
                    print(f"      ✗ Chunk {i} error: {e}")

        except Exception as e:
            print(f"      ✗ File error: {e}")
            failed_files += 1

    print(f"\n5. Creating vector index...")
    sys.stdout.flush()

    if len(all_vectors) == 0:
        print("   ❌ ERROR: No chunks created!")
        return

    # Create index
    vectors_array = np.array(all_vectors, dtype=np.float32)

    index_data = {
        "vectors": vectors_array,
        "metadata": all_metadata,
        "count": len(all_vectors),
        "dimension": vectors_array.shape[1],
    }

    # Upload index to S3
    print(f"\n6. Uploading index to S3...")
    sys.stdout.flush()
    s3.put_object(
        Bucket=S3_BUCKET,
        Key=INDEX_KEY,
        Body=pickle.dumps(index_data),
        ContentType="application/octet-stream",
    )

    print(f"\n{'=' * 80}")
    print(f"✅ REBUILD COMPLETE!")
    print(f"{'=' * 80}")
    print(f"Files processed: {processed_files}")
    print(f"Chunks created: {total_chunks}")
    print(f"Failed files: {failed_files}")
    print(f"Vector dimension: {vectors_array.shape[1]}")
    print(f"Index location: s3://{S3_BUCKET}/{INDEX_KEY}")
    print(f"\nNext: Restart backend to use the new index")


if __name__ == "__main__":
    main()
