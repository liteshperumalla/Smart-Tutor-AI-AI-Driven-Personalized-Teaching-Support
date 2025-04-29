import os
import streamlit as st
import logging
from html import escape
import smtplib
from email.message import EmailMessage
from llama_index.core import StorageContext, load_index_from_storage, get_response_synthesizer
from llama_index.core.schema import TextNode
from Tutor_chat import RAGQueryEngine
import re
import time
# New imports for in-app viewers
import base64
from io import BytesIO
from pathlib import Path
import streamlit.components.v1 as components
from pptx import Presentation
from docx import Document as DocxDoc
import mammoth
import json
import time
timestamp = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime())


# Import for speech recognition
import speech_recognition as sr

# Configure logging
logging.basicConfig(level=logging.ERROR, format='%(asctime)s - %(levelname)s - %(message)s')

# Disable Streamlit's file watcher
os.environ["STREAMLIT_SERVER_FILEWATCHER_TYPE"] = "none"

if "history" not in st.session_state:
    st.session_state.history = []  # list of tuples (role, message)
# ─── Chat Persistence Helpers ───────────────────────────────────────────────
PREV_CHAT_DIR = "previous_chats"
os.makedirs(PREV_CHAT_DIR, exist_ok=True)

def _safe_name(name: str) -> str:
    return re.sub(r'[^a-zA-Z0-9_-]', '_', name)

def _chat_path(name: str) -> str:
    return os.path.join(PREV_CHAT_DIR, f"{_safe_name(name)}.json")

def load_chat_sessions() -> dict[str, list]:
    sessions = {}
    for fname in os.listdir(PREV_CHAT_DIR):
        if not fname.endswith(".json"):
            continue
        path = os.path.join(PREV_CHAT_DIR, fname)
        try:
            with open(path, "r", encoding="utf-8") as f:
                sessions[fname[:-5]] = json.load(f)
        except:
            sessions[fname[:-5]] = []
    sessions.setdefault("Default", [])
    return sessions

def save_chat_session(name: str, history: list) -> None:
    formatted_history = []

    for entry in history:
        if isinstance(entry, (list, tuple)):
            if len(entry) == 3:
                # user: (role, message, timestamp)
                role, text, timestamp = entry
                formatted_history.append({
                    "role": role,
                    "content": text,
                    "timestamp": timestamp,
                    "sources": []
                })
            elif len(entry) == 4:
                # assistant: (role, message, sources, timestamp)
                role, text, sources, timestamp = entry
                formatted_history.append({
                    "role": role,
                    "content": text,
                    "timestamp": timestamp,
                    "sources": sources
                })
        elif isinstance(entry, dict):
            # Already correct
            formatted_history.append(entry)

    with open(_chat_path(name), "w", encoding="utf-8") as f:
        json.dump(formatted_history, f, ensure_ascii=False, indent=2)

# ---------- PAGE CONFIGURATION ----------
st.set_page_config(page_title="Smart AI Tutor", page_icon="🎓", layout="wide")
def _pdf_popup(path: str, page: int):
    """
    Opens a new browser window showing page=<page> of the PDF.
    Falls back to an alert if pop-ups are blocked.
    """
    data = Path(path).read_bytes()
    b64 = base64.b64encode(data).decode()
    url = f"data:application/pdf;base64,{b64}#page={page}"
    js = f"""
    <script type="text/javascript">
      var w = window.open("{url}", "_blank", "width=900,height=700");
      if (!w) {{
        alert("Pop-up blocked! Please allow pop-ups to view the PDF.");
      }}
    </script>
    """
    # height=0 so it doesn't take visual space
    st.markdown(js, unsafe_allow_html=True)
def _pptx_popup(path: str, slide_no: int):
    """
    Opens a popup window displaying the text of a single PPTX slide.
    """
    from pptx import Presentation
    prs   = Presentation(path)
    idx   = slide_no - 1
    if idx < 0 or idx >= len(prs.slides):
        st.error(f"Slide {slide_no} out of range.")
        return

    # Collect all text from the slide
    texts = []
    for shp in prs.slides[idx].shapes:
        if hasattr(shp, "text") and shp.text.strip():
            texts.append(shp.text.strip())
    slide_html = "<br>".join(texts).replace('"', "&quot;")

    js = f"""
    <script type="text/javascript">
      var w = window.open("", "_blank", "width=900,height=700");
      if (!w) {{ alert("Pop-up blocked! Please allow pop-ups."); }}
      else {{
        w.document.write(
          "<html><head><meta charset='utf-8'><title>Slide {slide_no}</title></head><body>" +
          "<h3>{os.path.basename(path)} — Slide {slide_no}</h3>" +
          "<div>{slide_html}</div>" +
          "</body></html>"
        );
      }}
    </script>
    """
    st.markdown(js, unsafe_allow_html=True)


def _docx_popup(path: str):
    """
    Opens a popup window displaying the full text of a DOCX.
    """
    import mammoth
    # Convert to HTML
    with open(path, "rb") as f:
        result = mammoth.convert_to_html(f)
    doc_html = result.value.replace('"', "&quot;")  # escape quotes

    js = f"""
    <script type="text/javascript">
      var w = window.open("", "_blank", "width=900,height=700");
      if (!w) {{ alert("Pop-up blocked! Please allow pop-ups."); }}
      else {{
        w.document.write(
          "<html><head><meta charset='utf-8'><title>{os.path.basename(path)}</title></head><body>" +
          "<h3>{os.path.basename(path)}</h3>" +
          `{doc_html}` +
          "</body></html>"
        );
      }}
    </script>
    """
    st.markdown(js, unsafe_allow_html=True)



# ---------- CSS STYLING ----------
LIGHT_MODE_CSS = """
<style>
  /* Headings & paragraphs black */
  h1, h2, h3, h4, p { color: black !important; }
  html, body, [class*="css"] { background-color: #f9fcff !important; color: black !important; }
  .main-title, .subtitle, .course-topics, .professor-title, .professor-name,
  .chat-title, .chat-subtitle { color: black !important; }
  .disclaimer {
    text-align: center;
    font-size: 0.8em;
    color: gray;
    margin-top: 20px; 
    padding-bottom: 10px;
  }
  /* Sidebar Button Styles */
  [data-testid="stSidebar"] button {
      background-color: transparent !important;
      color: inherit !important;
      border: none !important;
      text-align: left !important;
      padding: 10px 15px !important;
      border-radius: 8px !important;
      transition: background-color 0.2s;
  }
  [data-testid="stSidebar"] button:hover {
      background-color: #e0e0e0 !important;
      color: black !important;
  }
  [data-testid="stSidebar"] button:focus {
      background-color: #d6d6d6 !important;
  }
  [data-testid="stSidebar"] button:active {
      background-color: #cfcfcf !important;
  }
</style>
"""


DARK_MODE_CSS = """
<style>
  h1, h2, h3, h4, p { color: white !important; }
  [data-testid="stAppViewContainer"], [data-testid="stSidebar"],
  [data-testid="stHeader"], [data-testid="stToolbar"] {
      background-color: #000 !important;
      color: #fff !important;
  }
  .main-title, .subtitle, .course-topics,
  .professor-title, .professor-name,
  .chat-title, .chat-subtitle { color: #fff !important; }
  .disclaimer {
    text-align: center;
    font-size: 0.8em;
    color: gray;
    margin-top: 20px;
    padding-bottom: 10px;
  }
</style>
"""

BASE_CSS = """
<style>
  h1, h2, h3, h4, p { color: inherit !important; }
  .main-title { font-size: 3em; font-weight: 800; margin-bottom: 10px; }
  hr { border: none; height: 3px; background-color: #e0e0e0; margin-bottom: 25px; }
  .announcement-card { background: #fff3e0; border-left: 6px solid #fb8c00;
    border-radius: 8px; padding: 16px 20px; margin-bottom: 25px; max-height: 150px;
    overflow-y: auto;
  }
  .topics-card { background: #fff; border-radius: 8px; padding: 20px; margin-bottom: 25px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.08);
  }
  .professor-card { background: #fff; padding: 20px; border-radius: 8px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.1); text-align: center;
  }
  .professor-title { margin-top: 0; }
  .professor-name  { font-weight: 600; font-size: 1.1em; margin: 0; }
  .chat-container { display: inline-block; max-width: 70%; padding: 10px 15px;
    border-radius: 8px; margin-bottom: 10px; white-space: pre-wrap;
    word-wrap: break-word;
  }
  .user-chat      { background-color: #e3f2fd; margin-left: auto; }
  .assistant-chat { background-color: #f1f8e9; margin-right: auto; }
  .voice-button { margin-left: 10px; cursor: pointer; border: none;
    background-color: transparent; font-size: 20px;
  }
</style>
"""

# ---------- INDEX CONFIGURATION ----------
persist_dir = "./persisted_index"
storage_context = StorageContext.from_defaults(persist_dir=persist_dir)

# ---------- HELPERS & VIEWERS ----------
def is_greeting(msg: str) -> bool:
    greetings = ["hi", "hello", "hey", "good morning", "good evening",
                 "greetings", "thanks", "thank you", "bye", "goodbye"]
    return any(msg.lower().strip().startswith(g) for g in greetings)

def store_query(q: str):
    try:
        with open("query_log.txt", "a") as f:
            f.write(q + "\n")
    except Exception as e:
        logging.error(f"Error storing query: {e}")

def _show_pdf(path: str, page: int):
    pdf_bytes = Path(path).read_bytes()
    b64 = base64.b64encode(pdf_bytes).decode()
    src = f"data:application/pdf;base64,{b64}#page={page}"
    embed = f'<embed src="{src}" width="100%" height="800px" type="application/pdf">'
    components.html(embed, height=800)

def _show_pptx_text(path: str, slide_no: int):
    prs = Presentation(path)
    idx = slide_no - 1
    if idx < 0 or idx >= len(prs.slides):
        st.error(f"Slide {slide_no} is out of range.")
        return
    st.markdown(f"### {os.path.basename(path)} — Slide {slide_no}")
    for shape in prs.slides[idx].shapes:
        if hasattr(shape, "text") and shape.text.strip():
            st.write(shape.text.strip())

def _show_docx_text(path: str):
    doc = DocxDoc(path)
    st.markdown(f"### {os.path.basename(path)} (Word Document)")
    for para in doc.paragraphs:
        if para.text.strip():
            st.write(para.text.strip())

def generate_response_with_sources(q: str):
    """Returns (response_str, sources_list), where each source is a dict with file metadata."""
    if is_greeting(q):
        return "Hello! How can I assist you today?", []
    try:
        idx = load_index_from_storage(storage_context)
        retr = idx.as_retriever()
        synth = get_response_synthesizer(response_mode="compact")
        nodes = retr.retrieve(q)
        sources = []
        for n in nodes:
            fp = n.metadata.get("file_path")
            if not fp:
                continue

            try:
                chunk_text = n.get_text()
            except Exception:
                chunk_text = ""

            sources.append({
                "file_name": os.path.basename(fp),
                "file_path": fp,
                "page":      n.metadata.get("page_number"),
                "slide":     n.metadata.get("slide_number"),
                "chunk_text": chunk_text
            })

        qe = RAGQueryEngine(retriever=retr, response_synthesizer=synth)
        resp = qe.query(q)
        return str(resp), sources

    except Exception as e:
        err = f"⚠️ Error processing your query: {e}"
        st.error(err)
        logging.error(err)
        return err, []




def make_session_title(history):
    """
    Summarise the last few turns into a 3-word title.
    Handles history entries as (role,text) or {"role":…, "content":…}.
    """
    # Build a snippet of the last up to 6 messages
    snippet_lines = []
    for entry in history[-6:]:
        if isinstance(entry, (list, tuple)) and len(entry) >= 2:
            role, text = entry[0], entry[1]
        elif isinstance(entry, dict) and "role" in entry and "content" in entry:
            role, text = entry["role"], entry["content"]
        else:
            continue
        snippet_lines.append(f"{role}: {text}")

    snippet = "\n".join(snippet_lines)
    prompt = (
        "Summarize the following conversation in a concise, three-word title:\n"
        f"{snippet}\n"
        "Title (3 words):"
    )

    # Reuse your existing helper to query the LLM
    title, _ = generate_response_with_sources(prompt)

    # Extract up to 3 words, alphanumeric only
    words = re.findall(r"\w+", title)
    words = words[:3]
    if not words:
        return "Chat"
    return " ".join(words).title()

def sanitize_filename(name: str) -> str:
    """Lowercase, replace non-alnum with underscore, truncate to 50 chars."""
    safe = re.sub(r"[^0-9A-Za-z_-]", "_", name).lower()
    return safe[:50].rstrip("_")

def sidebar_content():
    with st.sidebar:
        st.markdown("### Navigation")
        for p, label in [
            ("home","Home"),("chat","Chat"),
            ("scheduleappointment","Schedule Appointment"),
            ("resources","Resources"),("about","About"),
        ]:
            if st.button(label):
                st.session_state.page = p
                st.rerun()

        st.markdown("<hr>", unsafe_allow_html=True)
        st.markdown("### Chats")

        # ─── init state ───────────────────────────────
        if "chat_sessions" not in st.session_state:
            st.session_state.chat_sessions = load_chat_sessions()
            st.session_state.current_chat = "Default"
        if "open_menu" not in st.session_state:
            st.session_state.open_menu = None

        # “New Chat” button
        if st.button("➕ New Chat"):
            name = f"Session {len(st.session_state.chat_sessions) + 1}"
            st.session_state.chat_sessions[name] = []
            st.session_state.current_chat = name
            save_chat_session(name, [])
            st.session_state.page = 'chat'
            st.session_state.open_menu = None
            st.rerun()

        # ─── build list: current first ─────────────────
        current = st.session_state.current_chat
        names = [current] + [n for n in st.session_state.chat_sessions if n != current]

        # ─── render each session with inline menu ──────
        for name in names:
            is_current = (name == current)

            # main row: session button + three‐dots
            with st.container():
                col1, col2 = st.columns([8,1], gap="small")
                with col1:
                    style = "background-color:#eef9ff; font-weight:bold;" if is_current else ""
                    if st.button(name, key=f"chat-{name}"):
                        st.session_state.current_chat = name
                        st.session_state.page = 'chat'
                        st.session_state.open_menu = None
                        st.rerun()
                    if is_current:
                        st.markdown(
                            f"<style>div[data-testid='stVerticalBlock'] > button[data-testid='stButton'] {{ {style} }}</style>",
                            unsafe_allow_html=True,
                        )
                with col2:
                    if st.button("⋮", key=f"menu-{name}"):
                        st.session_state.open_menu = None if st.session_state.open_menu == name else name

                # inline submenu if open
                if st.session_state.open_menu == name:
                    subcol1, subcol2, subcol3 = st.columns([6,1,1], gap="small")
                    with subcol1:
                        new_name = st.text_input(
                            "", value=name, key=f"rename-{name}",
                            placeholder="New name…", label_visibility="collapsed"
                        )
                    with subcol2:
                        if st.button("✔️", key=f"confirm-rename-{name}"):
                            sessions = st.session_state.chat_sessions
                            sessions[new_name] = sessions.pop(name)
                            if st.session_state.current_chat == name:
                                st.session_state.current_chat = new_name
                            save_chat_session(new_name, sessions[new_name])
                            st.session_state.open_menu = None
                            st.rerun()
                    with subcol3:
                        if st.button("🗑️", key=f"delete-{name}"):
                            # 1) Remove from state
                            sessions = st.session_state.chat_sessions
                            sessions.pop(name, None)
                            # 2) Delete file on disk
                            filename = sanitize_filename(name) + ".json"
                            filepath = os.path.join("previous_chats", filename)
                            try:
                                os.remove(filepath)
                            except FileNotFoundError:
                                pass
                            # 3) Pick a new current chat
                            remaining = list(sessions.keys())
                            st.session_state.current_chat = remaining[0] if remaining else None
                            # 4) Close menu and rerun
                            st.session_state.open_menu = None
                            st.rerun()

        st.markdown("<hr>", unsafe_allow_html=True)
        st.checkbox("Dark Mode", key="dark_mode")



# ---------- PAGES ----------
def home():
    st.markdown(
        "<div style='text-align:center;margin-bottom:10px;'>"
        "<h3>Welcome to Smart AI Tutor</h3>"
        "<div class='main-title'>INFO 5731 - Computational Methods</div>"
        "<div class='subtitle'>UNT | Fall 2025</div></div><hr>",
        unsafe_allow_html=True
    )
    c1, c2 = st.columns([3,1], gap="large")
    with c1:
        st.markdown(
            """
            <div class='announcement-card'>
                <strong>📢 Latest Announcements</strong><br><br>
                <p><strong>April 8, 2025:</strong> Assignment 3 released. Due by April 15.</p>
                <p style='color:#d50000;'><strong>[Reminder]</strong> Extra Credit Opportunity – Health Informatics Lecture Series: <em>Cybersecurity in Modern Healthcare</em> <strong>[April 9, 2025]</strong></p>
                <p><strong>April 5, 2025:</strong> Lecture notes updated.</p>
            </div>
            """,unsafe_allow_html=True)
        st.markdown(
    """
    <div class='topics-card'>
        <h4 class='course-topics' style='margin-top:0;'>Course Topics</h4>
        <ul class='course-topics'>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-1-lecture-materials?module_item_id=7518917">Intro to Python</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-2-lecture-materials?module_item_id=7518922">Python Basics 1</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-3-lecture-materials?module_item_id=7518928">Python Basics 2</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-4-lecture-materials?module_item_id=7518935">Web scraping using python</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-5-lecture-materials?module_item_id=7518942">Data Cleaning and Data Quality</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-6-lecture-materials?module_item_id=7518949">Feature Extraction</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-7-lecture-materials?module_item_id=7815101">Word Embedding and Transformer</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-10-lecture-materials-2?module_item_id=7518962">Topic Modeling</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-12-lecture-materials?module_item_id=7518972">Sentiment Analysis</a></li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-13-lecture-materials?module_item_id=7518976">Text Classification</li>
            <li><a href="https://unt.instructure.com/courses/117821/pages/week-14-lecture-materials?module_item_id=7518981">Generative AI in Natural Language Processing</a></li>
        </ul>
    </div>
    """,
    unsafe_allow_html=True)
    with c2:
        st.markdown(
    """
    <div class='professor-card'>
        <h4 class='professor-title'>Professor</h4>
        <p class='professor-name' style='font-size:20px; text-align:center;'>Dr. Haihua Chen</p>
        <p><a href='[https://www.linkedin.com/in/haihua-chen/](https://www.linkedin.com/in/haihua-chen/)' target='_blank'>LinkedIn</a><br>
        <a href='[https://scholar.google.com/citations?user=URmnWAQAAAAJ](https://scholar.google.com/citations?user=URmnWAQAAAAJ)' target='_blank'>Google Scholar</a></p>
    </div>
    """,
    unsafe_allow_html=True)
    st.markdown(
        "<hr><p class='disclaimer'>"  # Apply the disclaimer class
        "Disclaimer: The Smart AI Tutor may occasionally make mistakes. "
        "Please verify all important information independently."
        "</p>",
        unsafe_allow_html=True
    )

def chatbot():
    import re
    import time
    import base64
    from pathlib import Path
    from html import escape
    import streamlit as st

    # ─────────── Header ───────────
    st.markdown("<h1 style='text-align:center; margin-bottom:0.2em;'>Smart AI Tutor 🎓</h1>", unsafe_allow_html=True)
    st.markdown("<p class='chat-subtitle' style='text-align:center;'>Ask your questions:</p>", unsafe_allow_html=True)

    # ─────────── Bubble & CSS ───────────
    st.markdown("""
    <style>
    .user-bubble { background:#cce5ff; padding:10px; border-radius:15px; margin:10px; max-width:70%; text-align:justify; float:right; clear:both; }
    .assistant-bubble { background:#d4edda; padding:10px; border-radius:15px; margin:10px; max-width:70%; text-align:justify; float:left; clear:both; }
    .timestamp-user { font-size:0.7em; color:gray; text-align:right; margin-right:10px; margin-top:-5px; }
    .timestamp-assistant { font-size:0.7em; color:gray; text-align:left; margin-left:10px; margin-top:-5px; }
    pre { background:#000 !important; color:#fff !important; padding:10px !important; border-radius:8px !important; overflow-x:auto !important; }
    code { background:#000 !important; color:#fff !important; padding:2px 4px !important; border-radius:4px !important; }
    </style>
    """, unsafe_allow_html=True)

    # ─────────── State defaults ───────────
    if "chat_sessions" not in st.session_state:
        st.session_state.chat_sessions = {"Default": []}
    if "current_chat" not in st.session_state:
        st.session_state.current_chat = "Default"
    if "chat_input" not in st.session_state:
        st.session_state.chat_input = ""
    if "submitting" not in st.session_state:
        st.session_state.submitting = False

    history = st.session_state.chat_sessions[st.session_state.current_chat]

    # ─────────── Layout Containers ───────────
    chat_container = st.container()
    input_container = st.container()

    # ─────────── Submission callback ───────────
    def submit():
        query = st.session_state.chat_input.strip()
        if not query:
            return
        user_timestamp = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime())
        history.append(("user", query, user_timestamp))
        st.session_state.chat_input = ""
        st.session_state.submitting = True

    # ─────────── Render History ───────────
    with chat_container:
        seen_files = set()
        for msg_idx, msg in enumerate(history):
            role, text, sources, timestamp = None, None, [], ""

            if isinstance(msg, (list, tuple)):
                if len(msg) == 3:
                    role, text, timestamp = msg
                elif len(msg) == 4:
                    role, text, sources, timestamp = msg
            elif isinstance(msg, dict):
                role = msg.get("role")
                text = msg.get("content")
                sources = msg.get("sources", [])
                timestamp = msg.get("timestamp", "")

            if not role or not text:
                continue

            bubble_class = "user-bubble" if role == "user" else "assistant-bubble"
            html_content = re.sub(r"```(?:\w*\n)?(.*?)```", lambda m: f"<pre><code>{m.group(1)}</code></pre>", text, flags=re.DOTALL)
            st.markdown(f"<div class='{bubble_class}'>{html_content}</div>", unsafe_allow_html=True)

            if timestamp:
                timestamp_class = "timestamp-user" if role == "user" else "timestamp-assistant"
                st.markdown(f"<div class='{timestamp_class}'>{timestamp}</div>", unsafe_allow_html=True)

            if role == "assistant" and sources:
                for i, src in enumerate(sources):
                    fname = src.get("file_name", "file.pdf")
                    fpath = src.get("file_path", "")
                    chunk_text = src.get("chunk_text", "")
                    if not fpath or fpath in seen_files:
                        continue
                    seen_files.add(fpath)

                    try:
                        file_bytes = Path(fpath).read_bytes()
                    except Exception:
                        continue

                    key = f"dl-{st.session_state.current_chat}-{msg_idx}-{i}"
                    st.download_button(
                        label=f"📄 Download {fname}",
                        data=file_bytes,
                        file_name=fname,
                        mime="application/octet-stream",
                        key=key
                    )

                    if chunk_text:
                        escaped_text = escape(chunk_text).replace("\n", "<br>")
                        popup_html = f"""
                        <html><head><title>Highlighted Chunk</title></head><body>
                        <h3>{fname} — Matched Text</h3>
                        <div style='background-color:yellow; padding:10px; margin:10px; border-radius:8px;'>{escaped_text}</div>
                        </body></html>
                        """
                        highlight_button_html = f"""
                        <button onclick="var w=window.open('', '_blank', 'width=900,height=700');
                        if(w){{w.document.write(`{popup_html}`);}}else{{alert('Pop-up blocked! Please allow popups.');}}"
                        style="margin-top:10px; margin-bottom:10px; background-color:#ffd54f; border:none; border-radius:8px; padding:8px 16px; cursor:pointer;">
                        🔍 View Highlight
                        </button>
                        """
                        st.markdown(highlight_button_html, unsafe_allow_html=True)

        if st.session_state.get("submitting"):
            st.markdown("""
            <div class='assistant-bubble'>
             AI Tutor is thinking<span class="dot">.</span><span class="dot">.</span><span class="dot">.</span>
            </div>
            <style>
            .dot { animation: blink 1s infinite; }
            @keyframes blink { 0%{opacity:0;} 50%{opacity:1;} 100%{opacity:0;} }
            </style>
            """, unsafe_allow_html=True)

    # ─────────── Chat Input ───────────
    with input_container:
        st.text_input(
            "Ask",
            key="chat_input",
            placeholder="Ask me anything...",
            on_change=submit,
            label_visibility="collapsed"
        )
        
        if st.session_state.submitting:
            last_entry = history[-1]
            if isinstance(last_entry, (list, tuple)) and len(last_entry) >= 2:
                last_query = last_entry[1]
            elif isinstance(last_entry, dict) and "content" in last_entry:
                last_query = last_entry["content"]
            else:
                last_query = ""

            if last_query.strip():  # Avoid empty submissions
                reply, sources = generate_response_with_sources(last_query)

                assistant_timestamp = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime())
                history.append(("assistant", reply, sources, assistant_timestamp))

                # Session renaming logic
                old_name = st.session_state.current_chat
                sessions = st.session_state.chat_sessions
                new_name = make_session_title(history)

                if new_name not in sessions or new_name == old_name:
                    sessions[new_name] = sessions.pop(old_name)
                    st.session_state.current_chat = new_name

                safe = re.sub(r"[^0-9A-Za-z_-]", " ", new_name).lower()[:50].rstrip("_")
                save_chat_session(safe, sessions[st.session_state.current_chat])

            st.session_state.submitting = False
            st.rerun()


    # ─────────── Auto-scroll ─────────── 
    st.markdown("""
    <script>
    const chatContainer = window.parent.document.querySelector('.main');
    if (chatContainer) {
      chatContainer.scrollTo({ top: chatContainer.scrollHeight, behavior: 'smooth' });
    }
    </script>
    """, unsafe_allow_html=True)

    # ─────────── Disclaimer ───────────
    st.markdown("<p class='disclaimer'>The Smart AI Tutor may occasionally make mistakes. Please verify important information independently.</p>", unsafe_allow_html=True)

def appointment_page():
    import smtplib
    from email.message import EmailMessage
    import streamlit as st
    import logging

    st.markdown("<h2 style='text-align:center;'>📅 Schedule an Appointment</h2>", unsafe_allow_html=True)

    if 'form_submitted' not in st.session_state:
        st.session_state.form_submitted = False

    # ─────────── FORM ───────────
    with st.form('appt_form', clear_on_submit=True):
        name = st.text_input('👤 Your Name')
        email = st.text_input('📧 Your Email')
        date = st.date_input('📅 Preferred Date')
        t = st.time_input('⏰ Preferred Time')
        reason = st.text_area('📝 Reason (Optional)', placeholder="e.g., Discuss project queries, career advice...")

        submitted = st.form_submit_button('Submit Appointment Request')

    # ─────────── Submission Handling ───────────
    if submitted:
        if not name.strip():
            st.error("⚠️ Please enter your name.")
        elif '@' not in email or '.' not in email:
            st.error("⚠️ Please enter a valid email address.")
        else:
            try:
                # Compose Email
                msg = EmailMessage()
                msg['Subject'] = 'New Appointment Request - Smart AI Tutor'
                msg['From'] = st.secrets['email']['sender']
                msg['To'] = st.secrets['email']['recipient']
                body = f"""📅 New Appointment Request:\n\n
                Name: {name}\n
                Email: {email}\n
                Date: {date}\n
                Time: {t}\n
                Reason: {reason if reason else "N/A"}
                """
                msg.set_content(body)

                # Send Email
                with smtplib.SMTP_SSL(st.secrets['email']['smtp_server'], st.secrets['email']['smtp_port']) as server:
                    server.login(st.secrets['email']['username'], st.secrets['email']['password'])
                    server.send_message(msg)

                st.session_state.form_submitted = True

                # Stylish Confirmation
                st.success('🎉 Your appointment request has been submitted successfully!')
                st.balloons()

            except Exception as e:
                st.error(f'⚠️ Failed to send appointment request: {e}')
                logging.exception(e)

    # ─────────── Already Submitted ───────────
    if st.session_state.form_submitted and not submitted:
        st.success('✅ You have already submitted an appointment request.')
        st.info('📧 We will get back to you via email soon!')

    # ─────────── Footer Disclaimer ───────────
    st.markdown("""
    <hr>
    <p class='disclaimer' style='text-align:center; font-size:0.8em; color:gray;'>
    Disclaimer: Smart AI Tutor may occasionally make mistakes. Please verify critical information independently.
    </p>
    """, unsafe_allow_html=True)


def resources_page():
    """Displays the resources page."""
    st.markdown("<h2 style='text-align:center;'>📚 Resources</h2>", unsafe_allow_html=True)
    external_resources = {
        "Python Fundamentals": [
            {"title": "Official Python Documentation", "url": "[https://docs.python.org/3](https://docs.python.org/3)"},
            {"title": "PEP 8 Style Guide", "url": "[https://peps.python.org/pep-0008/](https://peps.python.org/pep-0008/)"},
            {"title": "Tutorialspoint: Python", "url": "[https://www.tutorialspoint.com/python/index.htm](https://www.tutorialspoint.com/python/index.htm)"},
            {"title": "Google's Python Class", "url": "[https://developers.google.com/edu/python](https://developers.google.com/edu/python)"},
            {"title": "Think Python Online", "url": "[http://greenteapress.com/thinkpython2/html/index.html](http://greenteapress.com/thinkpython2/html/index.html)"},
            {"title": "Complete Python Bootcamp (GitHub)", "url": "[https://github.com/Pierian-Data/Complete-Python-3-Bootcamp](https://github.com/Pierian-Data/Complete-Python-3-Bootcamp)"}
        ],
        "Development Tools & Environments": [
            {"title": "PyCharm (Education Access)", "url": "[https://www.jetbrains.com/community/education/#students](https://www.jetbrains.com/community/education/#students)"},
            {"title": "PyCharm Official Tutorial", "url": "[https://www.jetbrains.com/help/pycharm/creating-and-running-your-first-python-project.html](https://www.jetbrains.com/help/pycharm/creating-and-running-your-first-python-project.html)"},
            {"title": "Tutorialspoint: PyCharm", "url": "[https://www.tutorialspoint.com/pycharm/index.htm](https://www.tutorialspoint.com/pycharm/index.htm)"},
            {"title": "Google Colab Tutorial Notebook", "url": "[https://colab.research.google.com/github/ga642381/ML2021-Spring/blob/main/Colab/Google_Colab_Tutorial.ipynb](https://colab.research.google.com/github/ga642381/ML2021-Spring/blob/main/Colab/Google_Colab_Tutorial.ipynb)"},
            {"title": "Tutorialspoint: Google Colab", "url": "[https://www.tutorialspoint.com/google_colab/index.htm](https://www.tutorialspoint.com/google_colab/index.htm)"}
        ],
        "Python Tutorials & Courses": [
            {"title": "Data-Flair Python Tutorial for Beginners", "url": "[https://data-flair.training/blogs/python-tutorial](https://data-flair.training/blogs/python-tutorial)"},
            {"title": "Learn Python Programming (YouTube)", "url": "[https://www.youtube.com/watch?v=rfscVS0vtbw](https://www.youtube.com/watch?v=rfscVS0vtbw)"},
            {"title": "Coding Assistants with Generative AI", "url": "[https://github.com/steven2358/awesome-generative-ai?tab=readme-ov-file#coding](https://github.com/steven2358/awesome-generative-ai?tab=readme-ov-file#coding)"}
        ],
        "Natural Language Processing": [
            {"title": "NLTK 3.5 Documentation", "url": "[https://www.nltk.org/](https://www.nltk.org/)"},
            {"title": "Speech & Language Processing (Jurafsky & Martin)", "url": "[https://web.stanford.edu/~jurafsky/slp3](https://web.stanford.edu/~jurafsky/slp3)"},
            {"title": "NLP Is Fun!", "url": "[https://medium.com/@ageitgey/natural-language-processing-is-fun-9a0bff37854e](https://medium.com/@ageitgey/natural-language-processing-is-fun-9a0bff37854e)"},
            {"title": "Definitive Guide to NLP", "url": "[https://monkeylearn.com/blog/definitive-guide-natural-language-processing](https://monkeylearn.com/blog/definitive-guide-natural-language-processing)"},
            {"title": "Quora: How do I start with NLP?", "url": "[https://www.quora.com/How-do-I-start-with-Natural-Language-Processing](https://www.quora.com/How-do-I-start-with-Natural-Language-Processing)"},
            {"title": "Ask HN: Best Tools for Text Analysis", "url": "[https://news.ycombinator.com/item?id=9733883](https://news.ycombinator.com/item?id=9733883)"}
        ],
        "NLP Applications & Demos": [
            {"title": "Demo: Basic NLP (Colab)", "url": "[https://colab.research.google.com/drive/1JZnBoyjHy8QYgKEbZCGR9J0-P0qQycC0?usp=sharing](https://colab.research.google.com/drive/1JZnBoyjHy8QYgKEbZCGR9J0-P0qQycC0?usp=sharing)"},
            {"title": "Embedding Comparison (Colab)", "url": "[https://colab.research.google.com/drive/1t2EC7Aunf1qcsY4a50L4TFTXOM9SQB33?usp=sharing](https://colab.research.google.com/drive/1t2EC7Aunf1qcsY4a50L4TFTXOM9SQB33?usp=sharing)"},
            {"title": "In-Class Exercise 5 (Colab)", "url": "[https://colab.research.google.com/drive/17v-A-3qFlYYoHypoux7o2PvZ_e2RFsln?usp=sharing](https://colab.research.google.com/drive/17v-A-3qFlYYoHypoux7o2PvZ_e2RFsln?usp=sharing)"}
        ],
        "Sentiment Analysis & Reading": [
            {"title": "In-Depth Sentiment Analysis Series (Kaggle)", "url": "[https://www.kaggle.com/code/emirkocak/in-depth-series-sentiment-analysis-w-transformers](https://www.kaggle.com/code/emirkocak/in-depth-series-sentiment-analysis-w-transformers)"},
            {"title": "ChatGPT & Druid Sentiment App", "url": "[https://imply.io/blog/how-to-build-a-sentiment-analysis-application-with-chatgpt-and-druid](https://imply.io/blog/how-to-build-a-sentiment-analysis-application-with-chatgpt-and-druid)"},
            # Removed invalid PDF links
        ],
        "AI & Machine Learning": [
            {"title": "AI vs ML vs NLP vs DL (Quora)", "url": "[https://www.quora.com/What-is-the-difference-between-AI-Machine-Learning-NLP-and-Deep-Learning](https://www.quora.com/What-is-the-difference-between-AI-Machine-Learning-NLP-and-Deep-Learning)"},
            {"title": "Machine Learning Crash Course (Google)", "url": "[https://developers.google.com/machine-learning/crash-course/prereqs-and-prework](https://developers.google.com/machine-learning/crash-course/prereqs-and-prework)"},
            {"title": "Google AI Experiments", "url": "[https://experiments.withgoogle.com/collection/ai](https://experiments.withgoogle.com/collection/ai)"}
        ],
        "Word Embeddings & Language Models": [
            {"title": "What Are Word Embeddings?", "url": "[https://machinelearningmastery.com/what-are-word-embeddings/](https://machinelearningmastery.com/what-are-word-embeddings/)"},
            {"title": "Illustrated BERT & ELMo", "url": "[https://jalammar.github.io/illustrated-bert/](https://jalammar.github.io/illustrated-bert/)"},
            {"title": "ULMFiT Overview", "url": "[https://humboldt-wi.github.io/blog/research/information_systems_1819/group4_ulmfit](https://humboldt-wi.github.io/blog/research/information_systems_1819/group4_ulmfit)"},
            {"title": "BERT Model Survey", "url": "[https://towardsdatascience.com/a-review-of-bert-based-models-4ffdc0f15d58](https://towardsdatascience.com/a-review-of-bert-based-models-4ffdc0f15d58)"},
            {"title": "GPT-3 Reproduction Challenges", "url": "[https://jingfengyang.github.io/gpt](https://jingfengyang.github.io/gpt)"}
        ],
        "Text Mining & Topic Modeling": [
            {"title": "Text Mining 101", "url": "[https://tedunderwood.com/2012/08/14/where-to-start-with-text-mining](https://tedunderwood.com/2012/08/14/where-to-start-with-text-mining)"},
            {"title": "Topic Modeling with BERTopic", "url": "[https://github.com/MaartenGr/BERTopic](https://github.com/MaartenGr/BERTopic)"},
            {"title": "Gensim LDA Tutorial", "url": "[https://www.machinelearningplus.com/nlp/topic-modeling-gensim-python](https://www.machinelearningplus.com/nlp/topic-modeling-gensim-python)"}
        ],
        "Web Scraping & Automation": [
            {"title": "Twitter API v2 Quickstart", "url": "[https://developer.x.com/en/docs/tutorials/step-by-step-guide-to-making-your-first-request-to-the-twitter-api-v2](https://developer.x.com/en/docs/tutorials/step-by-step-guide-to-making-your-first-request-to-the-twitter-api-v2)"},
            {"title": "Web Scraping 101", "url": "[https://gregreda.com/2013/03/03/web-scraping-101-with-python](https://gregreda.com/2013/03/03/web-scraping-101-with-python)"},
            {"title": "ChatGPT Python Crawler Guide", "url": "[https://www.youtube.com/watch?v=B89Cf4pLNds](https://www.youtube.com/watch?v=B89Cf4pLNds)"},
            {"title": "GPT-4 Vision + Puppeteer", "url": "[https://www.youtube.com/watch?v=VeQR17k7fiU](https://www.youtube.com/watch?v=VeQR17k7fiU)"}
        ],
        "MLOps & Deployment": [
            {"title": "Model- to Data-Centric AI", "url": "[https://www.youtube.com/watch?v=06-AZXmwHjo](https://www.youtube.com/watch?v=06-AZXmwHjo)"}
        ],
        "Transformers & Conference Recordings": [
            {"title": "Transformers Session 1 (Zoom)", "url": "[https://unt.zoom.us/rec/share/L8-Zo_aX6QVsiOCcgZ84uNJRTNwybrNFv_WPZlZQWQr5d4NuQrBW5clscSxUbLtB.ryHtoLFwndhO6591](https://unt.zoom.us/rec/share/L8-Zo_aX6QVsiOCcgZ84uNJRTNwybrNFv_WPZlZQWQr5d4NuQrBW5clscSxUbLtB.ryHtoLFwndhO6591)"},
            {"title": "Transformers Session 2 (Zoom)", "url": "[https://unt.zoom.us/rec/share/GSTRUEL4xvoqJg-2_oS-7LmIhKl5jESNswaLbBaN84h7S59lcmGIUgmpQV9oa4WH.94JD-V1s9iGMjmuw](https://unt.zoom.us/rec/share/GSTRUEL4xvoqJg-2_oS-7LmIhKl5jESNswaLbBaN84h7S59lcmGIUgmpQV9oa4WH.94JD-V1s9iGMjmuw)"}
        ],
        "LangChain & LLM Utilities": [
            {"title": "What is LangChain?", "url": "[https://www.producthunt.com/stories/what-is-langchain-how-to-use](https://www.producthunt.com/stories/what-is-langchain-how-to-use)"},
            {"title": "LangChain Guides", "url": "[https://python.langchain.com/docs/use_cases/web_scraping](https://python.langchain.com/docs/use_cases/web_scraping)"},
            {"title": "KeyBERT Keyword Extraction", "url": "[https://maartengr.github.io/KeyBERT](https://maartengr.github.io/KeyBERT)"}
        ],
        "Neural Networks & GPT Projects": [
            {"title": "nanogpt Lecture (GitHub)", "url": "[https://github.com/karpathy/ng-video-lecture](https://github.com/karpathy/ng-video-lecture)"},
            {"title": "Zero to Hero Deep Learning", "url": "[https://karpathy.ai/zero-to-hero.html](https://karpathy.ai/zero-to-hero.html)"},
            {"title": "PicoGPT Repository", "url": "[https://github.com/jaymody/picoGPT/tree/29e78cc52b58ed2c1c483ffea2eb46ff6bdec785](https://github.com/jaymody/picoGPT/tree/29e78cc52b58ed2c1c483ffea2eb46ff6bdec785)"}
        ],
        "Generative AI Practices": [
            {"title": "GenAI Best Practices", "url": "[https://runawayhorse001.github.io/GenAI_Best_Practices/html/index.html](https://runawayhorse001.github.io/GenAI_Best_Practices/html/index.html)"},
            {"title": "ArXiv Preprint: GenAI", "url": "[https://arxiv.org/abs/2501.09223](https://arxiv.org/abs/2501.09223)"}
        ],
        "Social Sciences Python": [
            {"title": "Python Tutorials for Social Scientists", "url": "https://nealcaren.github.io/python-tutorials/"}
        ]
    }
    for topic, links in external_resources.items():
        with st.expander(f"📂 {topic}", expanded=False):
            for link in links:
                st.markdown(f"- [{link['title']}]({link['url']})", unsafe_allow_html=True)
    st.markdown(
        "<hr><p class='disclaimer'>"  # Apply the disclaimer class
        "Disclaimer: The Smart AI Tutor may occasionally make mistakes. "
        "Please verify all important information independently."
        "</p>",
        unsafe_allow_html=True
    )


def about_page():
    st.markdown("<h2 style='text-align:center;'>About Smart AI Tutor</h2>", unsafe_allow_html=True)
    st.markdown("<p>Smart AI Tutor is…</p>", unsafe_allow_html=True)
    st.markdown(
        "<hr><p class='disclaimer'>"  # Apply the disclaimer class
        "Disclaimer: The Smart AI Tutor may occasionally make mistakes. "
        "Please verify all important information independently."
        "</p>",
        unsafe_allow_html=True
    )


def main():
    if 'dark_mode' not in st.session_state:
        st.session_state.dark_mode = False
    if 'page' not in st.session_state:
        st.session_state.page = 'home'
    if 'chat_sessions' not in st.session_state:
        st.session_state.chat_sessions   = load_chat_sessions()
        st.session_state.current_chat    = 'Default'


    st.markdown(BASE_CSS, unsafe_allow_html=True)
    sidebar_content()
    st.markdown(DARK_MODE_CSS if st.session_state.dark_mode else LIGHT_MODE_CSS,
                unsafe_allow_html=True)

    if st.session_state.page=='home':
        home()
    elif st.session_state.page=='chat':
        chatbot()
    elif st.session_state.page=='scheduleappointment':
        appointment_page()
    elif st.session_state.page=='resources':
        resources_page()
    elif st.session_state.page=='about':
        about_page()

if __name__=='__main__':
    main()