"""
Simple rebuild: Process modules, generate Bedrock embeddings, upload to S3 with text
No complex dependencies - just core functionality
"""
import os
import sys
import json
import boto3
import pdfplumber
from pathlib import Path
from pptx import Presentation

sys.path.insert(0, '/app')
from backend.bedrock_embeddings import BedrockEmbeddings

print("="*80)
print("SIMPLE S3 REBUILD WITH BEDROCK + TEXT")
print("="*80)

# Config
S3_BUCKET = "smart-tutor-vector-bucket"
S3_REGION = "us-east-1"
MODULES_DIR = "/app/data/modules"
CHUNK_SIZE = 512

# Initialize
print("\n1. Initializing Bedrock...")
bedrock = BedrockEmbeddings()
print("   ✅ Bedrock ready")

s3 = boto3.client('s3', region_name=S3_REGION)

# Simple text extraction
def extract_text_from_pdf(pdf_path):
    """Extract text from PDF"""
    texts = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
    return "\n\n".join(texts)

def extract_text_from_pptx(pptx_path):
    """Extract text from PPTX"""
    texts = []
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            texts.append("\n".join(slide_text))
    return "\n\n".join(texts)

def simple_chunk(text, size=512):
    """Simple chunking by characters"""
    chunks = []
    for i in range(0, len(text), size):
        chunks.append(text[i:i+size])
    return chunks

# Process files
print(f"\n2. Processing files from {MODULES_DIR}...")
uploaded = 0

for root, dirs, files in os.walk(MODULES_DIR):
    for file in files:
        if file.endswith(('.pdf', '.pptx')):
            file_path = os.path.join(root, file)
            print(f"\n   Processing: {file}")

            try:
                # Extract text
                if file.endswith('.pdf'):
                    text = extract_text_from_pdf(file_path)
                else:
                    text = extract_text_from_pptx(file_path)

                if not text:
                    continue

                # Chunk
                chunks = simple_chunk(text, CHUNK_SIZE)
                print(f"      → {len(chunks)} chunks")

                # Upload each chunk
                for i, chunk_text in enumerate(chunks):
                    # Generate embedding
                    embedding = bedrock.get_text_embedding(chunk_text)

                    # Create chunk
                    chunk_id = f"{Path(file).stem}_chunk_{i:04d}"
                    chunk_data = {
                        'chunk_id': chunk_id,
                        'text': chunk_text,
                        'embedding': embedding,
                        'source_file': file,
                        'chunk_index': i
                    }

                    # Upload to S3
                    s3.put_object(
                        Bucket=S3_BUCKET,
                        Key=f"chunks/{chunk_id}.json",
                        Body=json.dumps(chunk_data),
                        ContentType='application/json'
                    )

                    uploaded += 1

                    if uploaded % 10 == 0:
                        print(f"      → Uploaded {uploaded} chunks...")

            except Exception as e:
                print(f"      ✗ Error: {e}")

print(f"\n{'='*80}")
print(f"✅ COMPLETE! Uploaded {uploaded} chunks to S3")
print(f"{'='*80}")
print(f"Location: s3://{S3_BUCKET}/chunks/")
