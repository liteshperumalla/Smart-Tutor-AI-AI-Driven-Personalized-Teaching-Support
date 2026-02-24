"""
Indexing Service — processes uploaded files into RAG-ready S3 chunks.

Flow per file:
  download from S3 → extract text → sentence-chunk → Bedrock embed
  → upload chunks/{resource_id}/… → invalidate prebuilt vector index

Progress is tracked in Redis (key: indexing:{resource_id}) with a 1-hour TTL.
Falls back to an in-memory dict when Redis is unavailable.
"""

from __future__ import annotations

import hashlib
import io
import json
import logging
import re
import threading
import time
import unicodedata
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Tuple

import boto3

from backend.config import config

logger = logging.getLogger(__name__)

# ── Status labels ─────────────────────────────────────────────────────────────
STATUS_QUEUED = "queued"
STATUS_EXTRACTING = "extracting"
STATUS_CHUNKING = "chunking"
STATUS_EMBEDDING = "embedding"
STATUS_UPLOADING = "uploading"
STATUS_COMPLETE = "complete"
STATUS_ERROR = "error"

# ── Chunking parameters (match process_modules_to_s3.py) ──────────────────────
CHUNK_SIZE = 512
CHUNK_OVERLAP = 102

# ── In-memory fallback for progress ──────────────────────────────────────────
_progress_store: dict = {}
_store_lock = threading.Lock()


# ── Text cleaning (copied verbatim from process_modules_to_s3.py) ─────────────
def _clean_text(text: str) -> str:
    text = unicodedata.normalize("NFKD", text)
    email_pat = r"([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})"
    emails = re.findall(email_pat, text)
    for i, e in enumerate(emails):
        text = text.replace(e, f"EMAIL_PLACEHOLDER_{i}")
    url_pat = r"http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+"
    urls = re.findall(url_pat, text)
    for i, u in enumerate(urls):
        text = text.replace(u, f"URL_PLACEHOLDER_{i}")
    text = re.sub(r"\.{5,}", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    for i, e in enumerate(emails):
        text = text.replace(f"EMAIL_PLACEHOLDER_{i}", e)
    for i, u in enumerate(urls):
        text = text.replace(f"URL_PLACEHOLDER_{i}", u)
    return text


class IndexingService:
    def __init__(self, redis_cache=None):
        self.redis = redis_cache
        self.s3 = boto3.client("s3", region_name=config.AWS_REGION)
        self.bucket = config.S3_DOCUMENTS_BUCKET

    # ── Progress helpers ───────────────────────────────────────────────────────

    def _set_progress(self, resource_id: str, data: dict):
        key = f"indexing:{resource_id}"
        if self.redis:
            try:
                self.redis.client.setex(key, 3600, json.dumps(data))
                return
            except Exception:
                pass
        with _store_lock:
            _progress_store[key] = data

    def get_status(self, resource_id: str) -> Optional[dict]:
        key = f"indexing:{resource_id}"
        if self.redis:
            try:
                val = self.redis.client.get(key)
                if val:
                    return json.loads(val)
            except Exception:
                pass
        with _store_lock:
            return _progress_store.get(key)

    # ── Public API ─────────────────────────────────────────────────────────────

    def start_indexing(
        self, resource_id: str, s3_key: str, filename: str, mime_type: str
    ):
        """Queue indexing in a background daemon thread."""
        self._set_progress(
            resource_id,
            {
                "status": STATUS_QUEUED,
                "progress_pct": 0,
                "chunks_created": 0,
                "total_chunks": None,
                "error": None,
                "started_at": datetime.utcnow().isoformat(),
                "completed_at": None,
            },
        )
        t = threading.Thread(
            target=self._run,
            args=(resource_id, s3_key, filename, mime_type),
            daemon=True,
        )
        t.start()

    # ── Background worker ──────────────────────────────────────────────────────

    def _run(self, resource_id: str, s3_key: str, filename: str, mime_type: str):
        started = datetime.utcnow().isoformat()

        def _prog(status, pct, created=0, total=None, error=None):
            self._set_progress(
                resource_id,
                {
                    "status": status,
                    "progress_pct": pct,
                    "chunks_created": created,
                    "total_chunks": total,
                    "error": error,
                    "started_at": started,
                    "completed_at": datetime.utcnow().isoformat() if status in (STATUS_COMPLETE, STATUS_ERROR) else None,
                },
            )

        try:
            # 1 — Download
            _prog(STATUS_EXTRACTING, 5)
            resp = self.s3.get_object(Bucket=self.bucket, Key=s3_key)
            file_bytes = resp["Body"].read()

            # 2 — Extract text
            ext = Path(filename).suffix.lower()
            pages = self._extract_text(file_bytes, filename, ext)
            if not pages:
                raise ValueError(f"No text extracted from '{filename}' (type: {ext})")

            # 3 — Chunk
            _prog(STATUS_CHUNKING, 30)
            chunks = self._chunk(pages)
            total = len(chunks)
            if total == 0:
                raise ValueError("Chunking produced zero chunks")

            # 4 — Embed + upload
            _prog(STATUS_EMBEDDING, 40, total=total)
            created = self._embed_and_upload(resource_id, filename, chunks, _prog)

            # 5 — Invalidate prebuilt index
            self._invalidate_index()

            _prog(STATUS_COMPLETE, 100, created=created, total=total)
            logger.info("Indexed %s → %d chunks (resource %s)", filename, created, resource_id)

        except Exception as exc:
            logger.error("Indexing failed for resource %s: %s", resource_id, exc)
            _prog(STATUS_ERROR, 0, error=str(exc))

    # ── Text extraction ────────────────────────────────────────────────────────

    def _extract_text(
        self, file_bytes: bytes, filename: str, ext: str
    ) -> List[Tuple[str, dict]]:
        pages: List[Tuple[str, dict]] = []

        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    raw = page.extract_text() or ""
                    if raw.strip():
                        pages.append((
                            _clean_text(raw),
                            {"file_name": filename, "page_number": page.page_number, "file_type": "pdf"},
                        ))

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            for idx, slide in enumerate(prs.slides, start=1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    pages.append((
                        _clean_text("\n".join(texts)),
                        {"file_name": filename, "slide_number": idx, "file_type": "pptx"},
                    ))

        elif ext == ".ipynb":
            nb = json.loads(file_bytes.decode("utf-8", errors="ignore"))
            parts = []
            for cell in nb.get("cells", []):
                if cell.get("cell_type") in ("markdown", "code"):
                    src = cell.get("source", [])
                    parts.append("".join(src) if isinstance(src, list) else src)
            joined = "\n\n".join(parts)
            if joined.strip():
                pages.append((_clean_text(joined), {"file_name": filename, "file_type": "notebook"}))

        elif ext in (".txt", ".md"):
            text = file_bytes.decode("utf-8", errors="ignore")
            if text.strip():
                pages.append((_clean_text(text), {"file_name": filename, "file_type": "text"}))

        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                if text.strip():
                    pages.append((_clean_text(text), {"file_name": filename, "file_type": "docx"}))
            except Exception as e:
                logger.warning("docx extraction failed: %s", e)

        return pages

    # ── Chunking (sentence-aware, no HF model needed) ─────────────────────────

    def _chunk(self, pages: List[Tuple[str, dict]]) -> List[Tuple[str, dict]]:
        chunks: List[Tuple[str, dict]] = []
        for text, meta in pages:
            if len(text) <= CHUNK_SIZE:
                if text.strip():
                    chunks.append((text, meta))
                continue
            sentences = re.split(r"(?<=[.!?])\s+", text)
            current: List[str] = []
            current_len = 0
            for sentence in sentences:
                sentence = sentence.strip()
                if not sentence:
                    continue
                if current_len + len(sentence) > CHUNK_SIZE and current:
                    chunk_text = " ".join(current)
                    if chunk_text.strip():
                        chunks.append((chunk_text, meta.copy()))
                    overlap = current[-1] if current else ""
                    current = [overlap, sentence] if overlap else [sentence]
                    current_len = len(overlap) + len(sentence)
                else:
                    current.append(sentence)
                    current_len += len(sentence)
            if current:
                chunk_text = " ".join(current)
                if chunk_text.strip():
                    chunks.append((chunk_text, meta.copy()))
        return chunks

    # ── Embed + upload ─────────────────────────────────────────────────────────

    def _embed_and_upload(
        self,
        resource_id: str,
        filename: str,
        chunks: List[Tuple[str, dict]],
        progress_cb,
    ) -> int:
        from backend.bedrock_embeddings import BedrockEmbeddings

        embeddings = BedrockEmbeddings(
            model_id=config.BEDROCK_EMBEDDING_MODEL_ID,
            region=config.AWS_REGION,
        )

        safe_name = re.sub(r"[^a-zA-Z0-9_.-]", "_", Path(filename).stem)[:30]
        base_path = f"resources/{resource_id}/{safe_name}"
        file_hash = hashlib.sha256(filename.encode()).hexdigest()[:8]
        total = len(chunks)
        created = 0

        for i, (text, _meta) in enumerate(chunks):
            try:
                embedding = embeddings.encode([text])[0]
            except Exception as e:
                logger.warning("Embedding failed for chunk %d: %s", i, e)
                continue

            chunk_id = f"{safe_name}_{file_hash}_chunk_{i:04d}"
            txt_key = f"chunks/{base_path}/chunk_{i:04d}.txt"
            vec_key = f"chunks/{base_path}/chunk_{i:04d}.vector.json"

            try:
                self.s3.put_object(
                    Bucket=self.bucket,
                    Key=txt_key,
                    Body=text.encode("utf-8"),
                    ContentType="text/plain",
                    Metadata={
                        "chunk-id": chunk_id,
                        "source-file": filename,
                        "chunk-index": str(i),
                    },
                )
                self.s3.put_object(
                    Bucket=self.bucket,
                    Key=vec_key,
                    Body=json.dumps({
                        "chunk_id": chunk_id,
                        "embedding": embedding,
                        "dimension": len(embedding),
                        "source_file": filename,
                        "chunk_index": i,
                    }).encode("utf-8"),
                    ContentType="application/json",
                )
                created += 1
            except Exception as e:
                logger.warning("S3 upload failed for chunk %d: %s", i, e)
                continue

            # Update progress (40% → 95%)
            pct = 40 + int((created / total) * 55)
            status = STATUS_UPLOADING if pct > 80 else STATUS_EMBEDDING
            progress_cb(status, pct, created=created, total=total)

            # Throttle slightly to avoid Bedrock rate limits
            time.sleep(0.05)

        return created

    # ── Invalidate prebuilt index ──────────────────────────────────────────────

    def _invalidate_index(self):
        """Remove the prebuilt vector index so the next query rebuilds from all chunks."""
        try:
            self.s3.delete_object(
                Bucket=self.bucket, Key="vector_index/s3_vector_index.pkl"
            )
            logger.info("Deleted prebuilt vector index — will rebuild on next query")
        except Exception as e:
            logger.warning("Could not delete prebuilt index (non-fatal): %s", e)


# ── Singleton ──────────────────────────────────────────────────────────────────
_indexing_service: Optional[IndexingService] = None


def get_indexing_service() -> IndexingService:
    global _indexing_service
    if _indexing_service is None:
        try:
            from backend.redis_cache import RedisCache
            rc = RedisCache(
                host=config.REDIS_HOST,
                port=config.REDIS_PORT,
                db=config.REDIS_DB,
                password=config.REDIS_PASSWORD,
                ssl=config.REDIS_SSL,
                max_connections=config.REDIS_MAX_CONNECTIONS,
            )
            _indexing_service = IndexingService(rc)
        except Exception as e:
            logger.warning("Redis unavailable for indexing service, using in-memory: %s", e)
            _indexing_service = IndexingService(None)
    return _indexing_service
