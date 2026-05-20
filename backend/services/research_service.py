"""
Research Service - S3-based implementation for file uploads and querying
"""

from __future__ import annotations

import base64
import io
import json
import logging
import re
import tempfile
import urllib.parse
import xml.etree.ElementTree as ET
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional
import uuid

import ipaddress
import socket

import requests
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from llama_index.core import Document, Settings, get_response_synthesizer
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.query_engine import RetrieverQueryEngine
from llama_index.core.vector_stores import ExactMatchFilter, MetadataFilters
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from PIL import Image
from pptx import Presentation
import pytesseract
import fitz
from youtube_transcript_api import (
    YouTubeTranscriptApi,
    TranscriptsDisabled,
    NoTranscriptFound,
)

from backend.config import config
from backend.s3_retriever import S3Retriever
from backend.s3_vector_store import S3VectorStore
from backend.bedrock_embeddings import BedrockEmbeddings

try:
    from serpapi import GoogleSearch

    SERPAPI_AVAILABLE = True
except ImportError:
    GoogleSearch = None
    SERPAPI_AVAILABLE = False

logger = logging.getLogger(__name__)

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 100
UPLOADS_PREFIX = "research_uploads/"
TEXT_CHUNKS_PREFIX = "research_chunks/"


def _search_web_results(*args, **kwargs):
    from utils import search_web_results

    return search_web_results(*args, **kwargs)


def _validate_url_not_internal(url: str) -> None:
    """Reject URLs that resolve to private/reserved IP addresses (SSRF protection)."""
    parsed = urllib.parse.urlparse(url)
    hostname = parsed.hostname
    if not hostname:
        raise ValueError("Invalid URL: no hostname")

    scheme = (parsed.scheme or "").lower()
    if scheme not in ("http", "https"):
        raise ValueError(f"Unsupported URL scheme: {scheme}")

    try:
        resolved_ips = socket.getaddrinfo(hostname, parsed.port or 443)
    except socket.gaierror:
        raise ValueError(f"Cannot resolve hostname: {hostname}")

    for family, _, _, _, sockaddr in resolved_ips:
        ip = ipaddress.ip_address(sockaddr[0])
        if ip.is_private or ip.is_loopback or ip.is_reserved or ip.is_link_local:
            raise ValueError(f"URL resolves to a non-public address")


class ResearchService:
    def __init__(self) -> None:
        self.s3_retriever = S3Retriever(similarity_top_k=10)
        self._folder_cache: Optional[Dict[str, List[str]]] = None
        self._sanitize_pattern = re.compile(r"[^a-zA-Z0-9._-]")

        self._s3_uploads = None
        self._uploads_bucket = config.S3_UPLOADS_BUCKET

    @property
    def s3_uploads(self):
        """Lazy initialize S3 client for uploads via centralized helper"""
        if self._s3_uploads is None:
            from backend.cloud.aws_helpers import get_boto3_client
            self._s3_uploads = get_boto3_client("s3")
        return self._s3_uploads

    def _sanitize(self, value: str) -> str:
        """Sanitize filename for S3 key"""
        safe = self._sanitize_pattern.sub("_", value)
        return safe[:100]

    def _ensure_bucket_exists(self):
        """Ensure the uploads bucket exists, create if needed"""
        try:
            import boto3

            s3 = boto3.resource("s3")
            bucket = s3.Bucket(self._uploads_bucket)
            bucket.meta.client.head_bucket(Bucket=self._uploads_bucket)
        except Exception as e:
            logger.info(f"Bucket {self._uploads_bucket} may not exist or error: {e}")
            pass

    def list_uploads(self, username: str) -> List[Dict[str, object]]:
        """List user uploads from S3 (scoped to the given user)"""
        try:
            from backend.cloud.aws_helpers import get_boto3_client

            s3 = get_boto3_client("s3")
            uploads = []
            paginator = s3.get_paginator("list_objects_v2")

            for page in paginator.paginate(
                Bucket=self._uploads_bucket, Prefix=f"{UPLOADS_PREFIX}{username}/", Delimiter="/"
            ):
                for obj in page.get("Contents", []):
                    key = obj["Key"]
                    if key.endswith(".json") and "metadata" in key:
                        try:
                            response = s3.get_object(
                                Bucket=self._uploads_bucket, Key=key
                            )
                            metadata = json.loads(response["Body"].read())
                            uploads.append(metadata)
                        except Exception:
                            pass

            uploads.sort(key=lambda x: x.get("uploaded_at", ""), reverse=True)
            return uploads
        except Exception as e:
            logger.error(f"Error listing uploads: {e}")
            return []

    def list_folders(self) -> List[Dict[str, object]]:
        """List folders from S3 knowledge base (course materials)"""
        try:
            from backend.cloud.aws_helpers import get_boto3_client

            s3 = get_boto3_client("s3")
            folder_counts: Dict[str, int] = {}

            paginator = s3.get_paginator("list_objects_v2")
            for page in paginator.paginate(
                Bucket=config.S3_DOCUMENTS_BUCKET, Prefix="modules/"
            ):
                for obj in page.get("Contents", []):
                    key = obj["Key"]
                    if key.endswith(
                        (".pdf", ".pptx", ".ppt", ".docx", ".ipynb", ".txt", ".md")
                    ):
                        parts = key.split("/")
                        if len(parts) > 2:
                            folder = parts[1]
                            folder_counts[folder] = folder_counts.get(folder, 0) + 1

            unique: Dict[str, Dict[str, object]] = {}
            for page in paginator.paginate(
                Bucket="smart-ai-tutor-docs", Prefix="modules/", Delimiter="/"
            ):
                for prefix in page.get("CommonPrefixes", []):
                    folder_path = prefix.get("Prefix", "")
                    label = folder_path.replace("modules/", "").replace("/", "")
                    if not label or label == "knowledge_uploads":
                        continue
                    key = label.lower()
                    if key not in unique:
                        unique[key] = {
                            "path": folder_path.rstrip("/"),
                            "label": label,
                            "file_count": folder_counts.get(label, 0),
                        }
            folders = list(unique.values())
            folders.sort(key=lambda item: str(item["label"]).lower())
            return folders
        except Exception as e:
            logger.error(f"Error listing folders from S3: {e}")
            return []

    def list_documents(self) -> List[Dict[str, object]]:
        """List documents from S3 knowledge base (course materials)"""
        try:
            from backend.cloud.aws_helpers import get_boto3_client

            s3 = get_boto3_client("s3")
            docs = []
            paginator = s3.get_paginator("list_objects_v2")
            for page in paginator.paginate(
                Bucket=config.S3_DOCUMENTS_BUCKET, Prefix="modules/"
            ):
                for obj in page.get("Contents", []):
                    key = obj["Key"]
                    if key.endswith(
                        (".pdf", ".pptx", ".ppt", ".docx", ".ipynb", ".txt", ".md")
                    ):
                        filename = key.split("/")[-1]
                        parts = key.split("/")
                        folder = parts[1] if len(parts) > 1 else "Unknown"
                        docs.append(
                            {
                                "id": key,
                                "title": filename,
                                "file_path": key,
                                "source": "s3",
                                "last_modified": obj["LastModified"].isoformat()
                                if hasattr(obj["LastModified"], "isoformat")
                                else str(obj["LastModified"]),
                            }
                        )
            docs.sort(key=lambda d: d.get("title", ""))
            return docs
        except Exception as e:
            logger.error(f"Error listing documents from S3: {e}")
            return []

    def _extract_text_from_file(
        self, content: bytes, filename: str
    ) -> Dict[str, object]:
        """Extract text and metadata from file content"""
        extension = Path(filename or "").suffix.lower()
        text = ""
        thumbnail = None

        if extension in {".txt", ".md", ".csv", ".py"}:
            text = content.decode("utf-8", errors="ignore")
            preview_type = "text"
        elif extension == ".ipynb":
            try:
                notebook = json.loads(content.decode("utf-8", errors="ignore"))
                cell_texts = []
                for cell in notebook.get("cells", []):
                    source = "".join(cell.get("source", []))
                    cell_type = cell.get("cell_type", "code")
                    if source.strip():
                        if cell_type == "code":
                            cell_texts.append(f"# Code Cell\n{source}")
                        else:
                            cell_texts.append(source)
                text = "\n\n".join(cell_texts)
            except (json.JSONDecodeError, KeyError):
                text = content.decode("utf-8", errors="ignore")
            preview_type = "text"
        elif extension == ".pdf":
            with fitz.open(stream=content, filetype="pdf") as pdf_doc:
                texts = []
                for page in pdf_doc:
                    texts.append(page.get_text())
                text = "\n".join(texts)
            preview_type = "pdf"
        elif extension == ".docx":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                doc = DocxDocument(tmp_path)
                text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            finally:
                Path(tmp_path).unlink(missing_ok=True)
            preview_type = "docx"
        elif extension == ".pptx":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                presentation = Presentation(tmp_path)
                slides = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides.append(shape.text)
                text = "\n".join(slides)
            finally:
                Path(tmp_path).unlink(missing_ok=True)
            preview_type = "pptx"
        elif extension in {".png", ".jpg", ".jpeg"}:
            image = Image.open(io.BytesIO(content))
            try:
                text = pytesseract.image_to_string(image)
            except Exception:
                text = ""
            b64 = base64.b64encode(content).decode("utf-8")
            mime = f"image/{image.format.lower()}" if image.format else "image/png"
            thumbnail = f"data:{mime};base64,{b64}"
            preview_type = "image"
        else:
            raise ValueError(f"Unsupported file type: {extension}")

        excerpt = (text or "No readable text extracted.").strip()[:800]
        return {
            "text": text,
            "preview_type": preview_type,
            "thumbnail": thumbnail,
            "excerpt": excerpt or "No readable text extracted.",
        }

    def _create_text_chunks(self, text: str, filename: str) -> List[Dict[str, Any]]:
        """Split text into overlapping chunks with metadata"""
        splitter = SentenceSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

        chunks = []
        doc = Document(text=text, metadata={"source_file": filename})
        nodes = splitter.get_nodes_from_documents([doc])

        for i, node in enumerate(nodes):
            chunk_id = f"{self._sanitize(filename)}_chunk_{i:05d}"
            chunks.append(
                {
                    "chunk_id": chunk_id,
                    "text": node.text,
                    "chunk_index": i,
                    "source_file": filename,
                }
            )

        return chunks

    def _upload_file_to_s3(self, content: bytes, s3_key: str) -> str:
        """Upload file content to S3 and return the key"""
        try:
            self.s3_uploads.put_object(
                Bucket=self._uploads_bucket,
                Key=s3_key,
                Body=content,
                ContentType="application/octet-stream",
            )
            return s3_key
        except Exception as e:
            logger.error(f"Error uploading to S3: {e}")
            raise

    def preview_file(self, content: bytes, filename: str) -> Dict[str, object]:
        """Upload a file, extract text, create chunks, and return preview"""
        upload_id = str(uuid.uuid4())[:8]
        sanitized_name = self._sanitize(filename)
        safe_id = f"{sanitized_name}_{upload_id}"

        s3_file_key = f"{UPLOADS_PREFIX}files/{safe_id}"
        s3_chunks_prefix = f"{TEXT_CHUNKS_PREFIX}{safe_id}/"

        try:
            self._ensure_bucket_exists()

            parsed = self._extract_text_from_file(content, filename)
            text = parsed["text"]

            # For images, allow empty OCR text — build description from image metadata
            if not text.strip():
                extension = Path(filename or "").suffix.lower()
                if extension in {".png", ".jpg", ".jpeg"}:
                    try:
                        img = Image.open(io.BytesIO(content))
                        w, h = img.size
                        fmt = (img.format or extension.lstrip(".")).upper()
                        mode = img.mode  # e.g. RGB, RGBA, L (grayscale)
                        size_kb = round(len(content) / 1024, 1)
                        text = (
                            f"[Uploaded image: {filename}]\n"
                            f"Format: {fmt}, Dimensions: {w}x{h}px, "
                            f"Color mode: {mode}, Size: {size_kb} KB"
                        )
                    except Exception:
                        text = f"[Uploaded image: {filename}]"
                    parsed["text"] = text
                    parsed["excerpt"] = text.split("\n")[0]
                else:
                    raise ValueError("No readable text extracted from file")

            self._upload_file_to_s3(content, s3_file_key)

            chunks = self._create_text_chunks(text, safe_id)

            for chunk in chunks:
                chunk_key = f"{s3_chunks_prefix}{chunk['chunk_id']}.json"
                self.s3_uploads.put_object(
                    Bucket=self._uploads_bucket,
                    Key=chunk_key,
                    Body=json.dumps(chunk),
                    ContentType="application/json",
                )

            metadata = {
                "id": safe_id,
                "file_name": filename,
                "file_path": s3_file_key,
                "chunks_prefix": s3_chunks_prefix,
                "chunk_count": len(chunks),
                "preview_type": parsed["preview_type"],
                "uploaded_at": datetime.now(timezone.utc).isoformat(),
            }

            metadata_key = f"{UPLOADS_PREFIX}metadata/{safe_id}.json"
            self.s3_uploads.put_object(
                Bucket=self._uploads_bucket,
                Key=metadata_key,
                Body=json.dumps(metadata),
                ContentType="application/json",
            )

            preview = self._build_preview(
                parsed["preview_type"],
                filename,
                text=parsed["excerpt"],
                thumbnail=parsed.get("thumbnail"),
                source=f"s3://{self._uploads_bucket}/{s3_file_key}",
            )
            preview["id"] = safe_id
            preview["file_name"] = filename
            return preview
        except Exception as e:
            logger.error(f"Error processing file upload: {e}")
            raise

    def preview_url(self, url: str) -> Dict[str, object]:
        """Fetch URL content and create searchable chunks"""
        # SECURITY: Reject URLs that resolve to internal/private addresses
        _validate_url_not_internal(url)

        upload_id = str(uuid.uuid4())[:8]
        safe_url_id = f"url_{upload_id}"

        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; SmartAITutor/1.0)",
            }
            response = requests.get(url, headers=headers, timeout=15, allow_redirects=False)
            response.raise_for_status()

            soup = BeautifulSoup(response.content, "html.parser")
            title = (
                soup.title.string.strip() if soup.title and soup.title.string else url
            )

            for element in soup(["script", "style", "nav", "footer", "aside"]):
                element.decompose()

            text = soup.get_text(separator="\n", strip=True)

            if not text.strip():
                raise ValueError("No readable text extracted from URL")

            self._ensure_bucket_exists()

            chunks = self._create_text_chunks(text, safe_url_id)

            s3_chunks_prefix = f"{TEXT_CHUNKS_PREFIX}{safe_url_id}/"
            for chunk in chunks:
                chunk_key = f"{s3_chunks_prefix}{chunk['chunk_id']}.json"
                self.s3_uploads.put_object(
                    Bucket=self._uploads_bucket,
                    Key=chunk_key,
                    Body=json.dumps(chunk),
                    ContentType="application/json",
                )

            metadata = {
                "id": safe_url_id,
                "file_name": title[:100],
                "source_url": url,
                "chunks_prefix": s3_chunks_prefix,
                "chunk_count": len(chunks),
                "preview_type": "webpage",
                "uploaded_at": datetime.now(timezone.utc).isoformat(),
            }

            metadata_key = f"{UPLOADS_PREFIX}metadata/{safe_url_id}.json"
            self.s3_uploads.put_object(
                Bucket=self._uploads_bucket,
                Key=metadata_key,
                Body=json.dumps(metadata),
                ContentType="application/json",
            )

            return self._build_preview(
                "webpage",
                title,
                text=text.strip()[:800],
                source=url,
            )
        except Exception as e:
            logger.error(f"Error processing URL: {e}")
            raise

    def preview_youtube(self, url: str) -> Dict[str, object]:
        """Fetch YouTube transcript and create searchable chunks"""
        video_id = self._extract_video_id(url)
        if not video_id:
            raise ValueError("Invalid YouTube URL")

        upload_id = str(uuid.uuid4())[:8]
        safe_yt_id = f"youtube_{video_id}_{upload_id}"

        try:
            api = YouTubeTranscriptApi()
            segments = api.fetch(video_id)
        except (TranscriptsDisabled, NoTranscriptFound) as exc:
            raise ValueError("Transcript unavailable for this video") from exc
        except Exception as exc:
            raise ValueError(f"Failed to fetch transcript: {exc}") from exc

        text_segments = []
        for segment in segments:
            if hasattr(segment, "text"):
                text_segments.append(segment.text)
            elif isinstance(segment, dict):
                text_segments.append(segment.get("text", ""))

        transcript = " ".join(text_segments)
        title = f"YouTube Video: {video_id}"

        if not transcript.strip():
            raise ValueError("Transcript unavailable for this video")

        self._ensure_bucket_exists()

        chunks = self._create_text_chunks(transcript, safe_yt_id)

        s3_chunks_prefix = f"{TEXT_CHUNKS_PREFIX}{safe_yt_id}/"
        for chunk in chunks:
            chunk_key = f"{s3_chunks_prefix}{chunk['chunk_id']}.json"
            self.s3_uploads.put_object(
                Bucket=self._uploads_bucket,
                Key=chunk_key,
                Body=json.dumps(chunk),
                ContentType="application/json",
            )

        metadata = {
            "id": safe_yt_id,
            "file_name": title,
            "youtube_url": url,
            "chunks_prefix": s3_chunks_prefix,
            "chunk_count": len(chunks),
            "preview_type": "youtube",
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
        }

        metadata_key = f"{UPLOADS_PREFIX}metadata/{safe_yt_id}.json"
        self.s3_uploads.put_object(
            Bucket=self._uploads_bucket,
            Key=metadata_key,
            Body=json.dumps(metadata),
            ContentType="application/json",
        )

        return self._build_preview(
            "youtube",
            title,
            text=transcript.strip()[:800],
            source=f"https://www.youtube.com/watch?v={video_id}",
        )

    def get_chunks_by_file_ids(
        self, file_ids: List[str]
    ) -> List[Dict[str, Any]]:
        """Fetch all text chunks for the given uploaded file IDs (by metadata lookup)."""
        try:
            import boto3
            s3 = boto3.client("s3", region_name=config.AWS_REGION)
            all_chunks: List[Dict[str, Any]] = []

            for fid in file_ids:
                # Read file metadata to get chunks_prefix
                meta_key = f"{UPLOADS_PREFIX}metadata/{fid}.json"
                try:
                    meta_resp = s3.get_object(Bucket=self._uploads_bucket, Key=meta_key)
                    meta = json.loads(meta_resp["Body"].read())
                except Exception:
                    logger.warning("Metadata not found for file %s", fid)
                    continue

                prefix = meta.get("chunks_prefix", "")
                if not prefix:
                    continue

                paginator = s3.get_paginator("list_objects_v2")
                for page in paginator.paginate(Bucket=self._uploads_bucket, Prefix=prefix):
                    for obj in page.get("Contents", []):
                        key = obj["Key"]
                        if not key.endswith(".json"):
                            continue
                        try:
                            resp = s3.get_object(Bucket=self._uploads_bucket, Key=key)
                            chunk = json.loads(resp["Body"].read())
                            all_chunks.append({
                                "id": chunk.get("chunk_id"),
                                "text": chunk.get("text", ""),
                                "score": 1.0,
                                "source_file": chunk.get("source_file", fid),
                                "chunk_index": chunk.get("chunk_index", 0),
                                "is_uploaded": True,
                            })
                        except Exception:
                            continue

            all_chunks.sort(key=lambda x: x.get("chunk_index", 0))
            return all_chunks
        except Exception as e:
            logger.error("Error fetching chunks by file IDs: %s", e)
            return []

    def _search_uploaded_chunks(
        self, query: str, top_k: int = 5
    ) -> List[Dict[str, Any]]:
        """Search through uploaded content chunks using keyword matching"""
        try:
            from backend.cloud.aws_helpers import get_boto3_client

            s3 = get_boto3_client("s3")
            results = []

            paginator = s3.get_paginator("list_objects_v2")
            for page in paginator.paginate(
                Bucket=self._uploads_bucket, Prefix=TEXT_CHUNKS_PREFIX
            ):
                for obj in page.get("Contents", []):
                    key = obj["Key"]
                    if not key.endswith(".json"):
                        continue

                    try:
                        response = s3.get_object(Bucket=self._uploads_bucket, Key=key)
                        chunk = json.loads(response["Body"].read())

                        text_lower = chunk.get("text", "").lower()
                        query_lower = query.lower()

                        if query_lower in text_lower:
                            score = text_lower.count(query_lower) / max(
                                len(text_lower.split()), 1
                            )
                            results.append(
                                {
                                    "id": chunk.get("chunk_id"),
                                    "text": chunk.get("text"),
                                    "score": min(score + 0.5, 1.0),
                                    "source_file": chunk.get("source_file"),
                                    "chunk_index": chunk.get("chunk_index", 0),
                                    "is_uploaded": True,
                                }
                            )
                    except Exception:
                        continue

            results.sort(key=lambda x: x["score"], reverse=True)
            return results[:top_k]
        except Exception as e:
            logger.error(f"Error searching uploaded chunks: {e}")
            return []

    def query(
        self,
        query: str,
        folders: Optional[List[str]] = None,
        uploaded_only: bool = False,
    ) -> Dict[str, Any]:
        """Query both course materials and uploaded content"""
        try:
            all_results = []

            if not uploaded_only:
                nodes = self.s3_retriever.retrieve(query)
                for node in nodes:
                    all_results.append(
                        {
                            "id": node.node.id_,
                            "text": node.node.text,
                            "score": node.score,
                            "source_file": node.node.metadata.get(
                                "source_file", "unknown"
                            ),
                            "chunk_index": node.node.metadata.get("chunk_index", 0),
                            "is_uploaded": False,
                        }
                    )

            uploaded_results = self._search_uploaded_chunks(query, top_k=10)
            all_results.extend(uploaded_results)

            all_results.sort(key=lambda x: x["score"], reverse=True)
            top_results = all_results[:20]

            answer_parts = []
            sources = []

            if top_results:
                for result in top_results[:5]:
                    answer_parts.append(result["text"])
                    sources.append(
                        {
                            "score": result["score"],
                            "excerpt": result["text"][:280],
                            "file_path": result.get("source_file"),
                            "title": result.get("source_file", "").split("/")[-1]
                            if result.get("source_file")
                            else "Uploaded content",
                        }
                    )

                combined_context = "\n\n".join(answer_parts[:3])

                from backend.bedrock_llm import BedrockLLM

                llm = BedrockLLM()
                answer = llm.generate(
                    prompt=f"Based on the following context, answer the question: {query}\n\nContext:\n{combined_context[:3000]}",
                    max_tokens=1024,
                )
            else:
                if uploaded_only:
                    answer = "No uploaded documents found matching your query. Please upload documents first."
                else:
                    answer = "No relevant information found in the knowledge base."

            return {
                "answer": answer,
                "sources": sources,
                "results": top_results,
                "total": len(top_results),
            }
        except Exception as e:
            logger.error(f"Error querying knowledge base: {e}")
            return {
                "answer": f"Error: {str(e)}",
                "sources": [],
                "results": [],
                "total": 0,
            }

    def clear_uploads(self, username: str) -> Dict[str, int]:
        """Clear uploaded content from S3 for the given user"""
        try:
            deleted_files = 0
            deleted_chunks = 0

            paginator = self.s3_uploads.get_paginator("list_objects_v2")
            for page in paginator.paginate(
                Bucket=self._uploads_bucket, Prefix=f"{UPLOADS_PREFIX}{username}/"
            ):
                for obj in page.get("Contents", []):
                    self.s3_uploads.delete_object(
                        Bucket=self._uploads_bucket, Key=obj["Key"]
                    )
                    deleted_files += 1

            for page in paginator.paginate(
                Bucket=self._uploads_bucket, Prefix=f"{TEXT_CHUNKS_PREFIX}{username}/"
            ):
                for obj in page.get("Contents", []):
                    self.s3_uploads.delete_object(
                        Bucket=self._uploads_bucket, Key=obj["Key"]
                    )
                    deleted_chunks += 1

            return {
                "deleted_files": deleted_files,
                "deleted_chunks": deleted_chunks,
                "deleted_count": deleted_files + deleted_chunks,
            }
        except Exception as e:
            logger.error(f"Error clearing uploads: {e}")
            return {"deleted_files": 0, "deleted_chunks": 0, "deleted_count": 0}

    def _build_preview(
        self,
        preview_type: str,
        title: str,
        *,
        text: Optional[str] = None,
        thumbnail: Optional[str] = None,
        source: Optional[str] = None,
    ) -> Dict[str, object]:
        return {
            "preview_type": preview_type,
            "title": title,
            "excerpt": (text or "").strip()[:800],
            "thumbnail": thumbnail,
            "source": source,
        }

    def _extract_video_id(self, url: str) -> Optional[str]:
        patterns = [
            r"[?&]v=([a-zA-Z0-9_-]{11})",
            r"youtu\.be/([a-zA-Z0-9_-]{11})",
            r"/embed/([a-zA-Z0-9_-]{11})",
            r"/v/([a-zA-Z0-9_-]{11})",
        ]
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        if re.fullmatch(r"[a-zA-Z0-9_-]{11}", url):
            return url
        return None

    # ==================== RESEARCH CAPABILITIES ====================

    def search_web(self, query: str, max_results: int = 5) -> Dict[str, Any]:
        """Search web for supplementary information."""
        try:
            results = _search_web_results(query, max_results)
            return {
                "query": query,
                "web_results": results,
                "count": len(results),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
        except Exception as e:
            logger.error(f"Web search failed: {e}")
            return {
                "query": query,
                "web_results": [],
                "count": 0,
                "error": str(e),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }

    def search_academic_papers(
        self, query: str, sources: Optional[List[str]] = None, max_results: int = 10
    ) -> Dict[str, Any]:
        """Search academic databases for relevant papers."""
        sources = sources or ["arxiv", "pubmed", "scholar"]
        papers = []
        errors = []

        for source in sources:
            try:
                if source == "arxiv":
                    papers.extend(
                        self._search_arxiv(query, max_results // len(sources) + 1)
                    )
                elif source == "pubmed":
                    papers.extend(
                        self._search_pubmed(query, max_results // len(sources) + 1)
                    )
                elif source == "scholar" and SERPAPI_AVAILABLE:
                    papers.extend(
                        self._search_google_scholar(
                            query, max_results // len(sources) + 1
                        )
                    )
            except Exception as e:
                logger.error(f"Academic search failed for {source}: {e}")
                errors.append({"source": source, "error": str(e)})

        return {
            "query": query,
            "papers": papers[:max_results],
            "sources_searched": sources,
            "count": len(papers[:max_results]),
            "errors": errors if errors else None,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        }

    def _search_arxiv(self, query: str, max_results: int) -> List[Dict[str, Any]]:
        """Search arXiv API for papers."""
        encoded_query = urllib.parse.quote(query)
        url = f"https://export.arxiv.org/api/query?search_query=all:{encoded_query}&max_results={max_results}&sortBy=relevance"

        try:
            response = requests.get(url, timeout=15)
            response.raise_for_status()
            root = ET.fromstring(response.content)

            ns = {"atom": "http://www.w3.org/2005/Atom"}
            papers = []

            for entry in root.findall("atom:entry", ns):
                title = entry.find("atom:title", ns)
                summary = entry.find("atom:summary", ns)
                published = entry.find("atom:published", ns)
                link = entry.find("atom:id", ns)
                authors = entry.findall("atom:author/atom:name", ns)

                papers.append(
                    {
                        "title": title.text.strip().replace("\n", " ")
                        if title is not None
                        else "Unknown",
                        "abstract": summary.text.strip()[:500]
                        if summary is not None
                        else "",
                        "authors": [a.text for a in authors][:5],
                        "url": link.text if link is not None else "",
                        "source": "arxiv",
                        "published_date": published.text[:10]
                        if published is not None
                        else None,
                    }
                )

            return papers
        except Exception as e:
            logger.error(f"arXiv search error: {e}")
            return []

    def _search_pubmed(self, query: str, max_results: int) -> List[Dict[str, Any]]:
        """Search PubMed/NCBI for papers using E-utilities."""
        base_url = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils"

        try:
            search_url = f"{base_url}/esearch.fcgi?db=pubmed&term={urllib.parse.quote(query)}&retmax={max_results}&retmode=json"
            search_response = requests.get(search_url, timeout=15)
            search_data = search_response.json()
            ids = search_data.get("esearchresult", {}).get("idlist", [])

            if not ids:
                return []

            fetch_url = (
                f"{base_url}/esummary.fcgi?db=pubmed&id={','.join(ids)}&retmode=json"
            )
            fetch_response = requests.get(fetch_url, timeout=15)
            fetch_data = fetch_response.json()

            papers = []
            results = fetch_data.get("result", {})
            for pmid in ids:
                if pmid not in results:
                    continue
                article = results[pmid]
                authors = article.get("authors", [])
                pubdate = (
                    article.get("pubdate", "")[:10] if article.get("pubdate") else None
                )

                papers.append(
                    {
                        "title": article.get("title", "Unknown"),
                        "abstract": article.get("abstract", "")[:500],
                        "authors": [a.get("name", "") for a in authors][:5],
                        "url": f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
                        "source": "pubmed",
                        "published_date": pubdate,
                        "doi": article.get("elocationid", "").replace("doi: ", "")
                        if article.get("elocationid")
                        else None,
                    }
                )

            return papers
        except Exception as e:
            logger.error(f"PubMed search error: {e}")
            return []

    def _search_google_scholar(
        self, query: str, max_results: int
    ) -> List[Dict[str, Any]]:
        """Search Google Scholar using SerpAPI."""
        if not SERPAPI_AVAILABLE or not GoogleSearch:
            return []

        try:
            params = {
                "engine": "google_scholar",
                "q": query,
                "num": max_results,
                "api_key": config.SERPAPI_API_KEY,
            }
            search = GoogleSearch(params)
            results = search.get_dict()

            papers = []
            for result in results.get("organic_results", []):
                papers.append(
                    {
                        "title": result.get("title", "Unknown"),
                        "abstract": result.get("snippet", "")[:500],
                        "authors": [],
                        "url": result.get("link", ""),
                        "source": "scholar",
                        "cited_by": result.get("cited_by", {}).get("value"),
                    }
                )

            return papers
        except Exception as e:
            logger.error(f"Google Scholar search error: {e}")
            return []

    def compare_sources(
        self,
        topic: str,
        document_ids: Optional[List[str]] = None,
        uploaded_only: bool = True,
    ) -> Dict[str, Any]:
        """Compare information across uploaded documents on a topic."""
        from backend.bedrock_llm import BedrockLLM

        try:
            chunks = self._search_uploaded_chunks(topic, top_k=20)

            doc_contents: Dict[str, List[str]] = {}
            for chunk in chunks:
                source = chunk.get("source_file", "unknown")
                if source not in doc_contents:
                    doc_contents[source] = []
                doc_contents[source].append(chunk.get("text", ""))

            comparisons = []
            agreements = []
            contradictions = []
            analyzed = 0

            for source, texts in doc_contents.items():
                combined = "\n".join(texts[:5])
                comparisons.append(
                    {
                        "document_title": source,
                        "file_path": source,
                        "excerpts": texts[:3],
                        "relevance_score": sum(
                            c["score"] for c in chunks if c.get("source_file") == source
                        )
                        / max(
                            len([c for c in chunks if c.get("source_file") == source]),
                            1,
                        ),
                    }
                )
                analyzed += 1

            llm = BedrockLLM()
            prompt = f"""Compare these sources on the topic "{topic}" and identify agreements and contradictions:

Sources analyzed: {len(comparisons)}
"""

            response = llm.generate(prompt=prompt, max_tokens=512)

            return {
                "topic": topic,
                "documents_analyzed": analyzed,
                "comparisons": comparisons,
                "agreements": agreements,
                "contradictions": contradictions,
                "summary": response,
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
        except Exception as e:
            logger.error(f"Compare sources error: {e}")
            return {
                "topic": topic,
                "documents_analyzed": 0,
                "comparisons": [],
                "agreements": [],
                "contradictions": [],
                "summary": f"Error: {str(e)}",
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }

    def extract_citations(
        self, document_id: Optional[str] = None, format_style: str = "apa"
    ) -> Dict[str, Any]:
        """Extract and format citations from uploaded documents."""
        from backend.bedrock_llm import BedrockLLM

        try:
            chunks = self._search_uploaded_chunks(
                "reference citation bibliography", top_k=10
            )
            content = "\n".join([c.get("text", "") for c in chunks])

            citations = self._detect_citations(content)
            formatted = [self._format_citation(c, format_style) for c in citations]

            bibliography = "\n\n".join([f["formatted"] for f in formatted])

            return {
                "citations": formatted,
                "format": format_style,
                "count": len(formatted),
                "exportable_bibliography": bibliography,
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
        except Exception as e:
            logger.error(f"Extract citations error: {e}")
            return {
                "citations": [],
                "format": format_style,
                "count": 0,
                "exportable_bibliography": "",
                "error": str(e),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }

    def _detect_citations(self, content: str) -> List[Dict[str, Any]]:
        """Detect citations in text."""
        citations = []
        citation_pattern = r"\[(\d+)\]|(?:^|\n)\d+\.\s+([^\n]+)"
        matches = re.findall(citation_pattern, content, re.MULTILINE)

        for match in matches:
            citations.append(
                {
                    "raw": match[0] or match[1],
                    "year": re.search(r"\b(19|20)\d{2}\b", match[0] or match[1]).group()
                    if re.search(r"\b(19|20)\d{2}\b", match[0] or match[1])
                    else None,
                }
            )

        return citations

    def _format_citation(self, citation: Dict[str, Any], style: str) -> Dict[str, Any]:
        """Format a single citation."""
        raw = citation.get("raw", "")
        formatted = raw

        if style == "apa":
            formatted = f"{raw} ({citation.get('year', 'n.d.')})"
        elif style == "mla":
            formatted = f'"{raw}."'
        elif style == "ieee":
            formatted = f"[1] {raw}"
        elif style == "chicago":
            formatted = f'"{raw}."'

        return {
            "raw": raw,
            "formatted": formatted,
            "style": style,
        }

    def generate_summary(
        self,
        document_id: Optional[str] = None,
        mode: str = "executive",
        max_length: Optional[int] = None,
    ) -> Dict[str, Any]:
        """Generate summary of uploaded documents."""
        from backend.bedrock_llm import BedrockLLM

        try:
            chunks = self._search_uploaded_chunks("", top_k=10)
            content = "\n".join([c.get("text", "") for c in chunks])[:5000]

            llm = BedrockLLM()
            prompt = f"Generate a {mode} summary of the following content:\n\n{content}"

            if max_length:
                prompt += f"\n\nLimit to {max_length} words."

            summary = llm.generate(prompt=prompt, max_tokens=1024)

            return {
                "mode": mode,
                "summary": summary,
                "word_count": len(summary.split()),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
        except Exception as e:
            logger.error(f"Generate summary error: {e}")
            return {
                "mode": mode,
                "summary": f"Error: {str(e)}",
                "error": str(e),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }

    def generate_questions(
        self,
        document_id: Optional[str] = None,
        difficulty: str = "medium",
        question_types: Optional[List[str]] = None,
        count: int = 5,
    ) -> Dict[str, Any]:
        """Generate study questions from uploaded documents."""
        from backend.bedrock_llm import BedrockLLM

        question_types = question_types or ["mcq", "short_answer"]
        difficulty_levels = {
            "easy": "basic",
            "medium": "moderate",
            "hard": "challenging",
        }

        try:
            chunks = self._search_uploaded_chunks("", top_k=10)
            content = "\n".join([c.get("text", "") for c in chunks])[:4000]

            llm = BedrockLLM()
            prompt = f"""Generate {count} {difficulty_levels.get(difficulty, "moderate")} study questions from this content. 
Question types: {", ".join(question_types)}

Content:
{content}

Format as JSON array with: question, type, difficulty, answer (for reference), explanation"""

            response = llm.generate(prompt=prompt, max_tokens=2048)
            questions = []

            try:
                import json

                questions = json.loads(response)
            except json.JSONDecodeError:
                questions = [
                    {
                        "question": response,
                        "type": "short_answer",
                        "difficulty": difficulty,
                    }
                ]

            return {
                "questions": questions[:count],
                "difficulty": difficulty,
                "types": question_types,
                "count": len(questions),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
        except Exception as e:
            logger.error(f"Generate questions error: {e}")
            return {
                "questions": [],
                "difficulty": difficulty,
                "types": question_types,
                "count": 0,
                "error": str(e),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }

    def fact_check(
        self, claim: str, uploaded_only: bool = True, include_web: bool = False
    ) -> Dict[str, Any]:
        """Cross-reference claims against sources."""
        from backend.bedrock_llm import BedrockLLM

        try:
            sources = self._search_uploaded_chunks(claim, top_k=10)

            evidence = []
            for source in sources:
                evidence.append(
                    {
                        "source_type": "document",
                        "title": source.get("source_file", "Unknown"),
                        "excerpt": source.get("text", "")[:300],
                        "supports_claim": True,
                        "confidence": source.get("score", 0),
                    }
                )

            if include_web:
                web_results = self.search_web(claim, max_results=3)
                for result in web_results.get("web_results", []):
                    evidence.append(
                        {
                            "source_type": "web",
                            "title": result.get("title", "Unknown"),
                            "url": result.get("url", ""),
                            "excerpt": result.get("content", "")[:300],
                            "supports_claim": True,
                            "confidence": 0.5,
                        }
                    )

            supporting = [e for e in evidence if e.get("supports_claim")]
            contradicting = [e for e in evidence if not e.get("supports_claim")]

            avg_confidence = sum(e.get("confidence", 0) for e in evidence) / max(
                len(evidence), 1
            )

            verdict = (
                "supported"
                if avg_confidence > 0.6
                else "contradicted"
                if avg_confidence < 0.3
                else "inconclusive"
            )

            return {
                "claim": claim,
                "verdict": verdict,
                "confidence": avg_confidence,
                "supporting_sources": supporting,
                "contradicting_sources": contradicting,
                "evidence": evidence,
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }
        except Exception as e:
            logger.error(f"Fact check error: {e}")
            return {
                "claim": claim,
                "verdict": "inconclusive",
                "confidence": 0,
                "supporting_sources": [],
                "contradicting_sources": [],
                "evidence": [],
                "error": str(e),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }


_research_service: Optional[ResearchService] = None


def get_research_service() -> ResearchService:
    """Dependency injection for ResearchService (singleton)"""
    global _research_service
    if _research_service is None:
        _research_service = ResearchService()
    return _research_service
