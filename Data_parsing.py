import os
import torch
import re
import unicodedata
import pdfplumber
from pathlib import Path
from pptx import Presentation
from transformers import AutoModel, AutoTokenizer
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings, get_response_synthesizer
from llama_index.core.ingestion import IngestionPipeline
from llama_index.core.node_parser import SentenceSplitter
from llama_index.vector_stores.chroma import ChromaVectorStore
from chromadb import PersistentClient
from llama_index.core.retrievers import BaseRetriever
from llama_index.core.query_engine import CustomQueryEngine
from llama_index.core.response_synthesizers import BaseSynthesizer
from llama_index.llms.ollama import Ollama
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document

# ------------------------------
# CUSTOM PPTX READER
# ------------------------------
class PPTXTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        prs = Presentation(file_path)
        docs: list[Document] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw = shape.text.strip()
                    if raw:
                        text_runs.append(preprocess_text(file_path, raw))
            slide_text = "\n".join(text_runs)
            if slide_text:
                docs.append(Document(
                    text=slide_text,
                    metadata={
                        "file_path": file_path,
                        "slide_number": slide_idx
                    }
                ))
        return docs

# ------------------------------
# CUSTOM PDF READER
# ------------------------------
class PDFTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        docs: list[Document] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                raw = page.extract_text()
                if raw:
                    cleaned = preprocess_text(file_path, raw)
                    docs.append(Document(
                        text=cleaned,
                        metadata={
                            "file_path": file_path,
                            "page_number": page.page_number
                        }
                    ))
        return docs
# ------------------------------
# MODEL & LLM SETTINGS
# ------------------------------
try:
    model_name = "sentence-transformers/all-MiniLM-L6-v2"
    Settings.embed_model = HuggingFaceEmbedding(model_name=model_name)
    print(f"✅ Model {model_name} loaded successfully.")
except Exception as e:
    print(f"❌ Error loading embedding model: {e}")
    exit()

Settings.llm = Ollama(model="llama3.1:latest", request_timeout=120.0)

# ------------------------------
# TEXT PREPROCESSING FUNCTION
# ------------------------------
def preprocess_text(file_path, text):
    code_extensions = {".py", ".java", ".cpp", ".js", ".c", ".cs", ".html", ".css", ".php", ".rb"}
    text_extensions = {".pdf", ".docx", ".pptx", ".txt"}

    ext = os.path.splitext(file_path)[-1].lower()

    if ext in text_extensions:
        text = clean_text(text)
    return text

def clean_text(text):
    text = unicodedata.normalize("NFKD", text)

    # Preserve email addresses
    email_pattern = r'([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})'
    emails = re.findall(email_pattern, text)
    for i, email in enumerate(emails):
        text = text.replace(email, f'EMAIL_PLACEHOLDER_{i}')

    # Preserve URLs
    url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
    urls = re.findall(url_pattern, text)
    for i, url in enumerate(urls):
        text = text.replace(url, f'URL_PLACEHOLDER_{i}')

    # Remove table of contents artifacts (long sequences of dots)
    text = re.sub(r'\.{5,}', ' ', text)  # Remove sequences of 5 or more dots

    # Remove excessive spaces and unnecessary formatting
    text = re.sub(r'\s+', ' ', text).strip()  

    # Restore emails and URLs
    for i, email in enumerate(emails):
        text = text.replace(f'EMAIL_PLACEHOLDER_{i}', email)
    for i, url in enumerate(urls):
        text = text.replace(f'URL_PLACEHOLDER_{i}', url)

    return text

# ------------------------------
# LOAD DOCUMENTS & PREPROCESS
# ------------------------------
doc_path = "/Users/liteshperumalla/Desktop/Files/masters/Smart AI Tutor/Modules/"

try:
    reader = SimpleDirectoryReader(
        input_dir=doc_path,
        required_exts=['.pptx', '.ipynb', '.docx', '.csv', '.jpeg', '.pdf', '.png', '.py'],
        file_extractor={".pptx": PPTXTextOnlyReader(), ".pdf": PDFTextOnlyReader()}, 
        recursive=True
    )
    docs = reader.load_data()

    if not docs:
        print("❌ No documents found! Check the path and file extensions.")
        exit()

    print(f"✅ Loaded {len(docs)} docs")

except Exception as e:
    print(f"❌ Error loading documents: {e}")
    exit()

# ------------------------------
# PROCESS DOCUMENTS & CHUNKING
# ------------------------------
processed_docs = []
CHUNK_SIZE = 512
CHUNK_OVERLAP = 50

def chunk_text(text, chunk_size=512, overlap=50):
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end]
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

# Preprocess and chunk the documents
for doc in docs:
    # get page or slide number from metadata, if any
    slide_no = doc.metadata.get("slide_number")
    page_no  = doc.metadata.get("page_number")

    # chunk the text
    chunks = chunk_text(doc.get_content(), CHUNK_SIZE, CHUNK_OVERLAP)
    for chunk in chunks:
        md = {
            "file_name":  os.path.basename(doc.metadata["file_path"]),
            "file_path":  doc.metadata["file_path"],
            "folder_name": os.path.basename(os.path.dirname(doc.metadata["file_path"])),
            "num_tokens": len(chunk.split()),
            "num_chars":  len(chunk),
        }
        if slide_no is not None:
            md["slide_number"] = slide_no
        if page_no is not None:
            md["page_number"] = page_no

        processed_docs.append({
            "doc_id":   doc.doc_id,
            "text":     chunk,
            "metadata": md,
            "category": "<category>"
        })

# ------------------------------
# CONVERT PROCESSED DOCS TO DOCUMENT OBJECTS
# ------------------------------
document_objects = []
for doc in processed_docs:
    document_objects.append(Document(text=doc["text"], metadata=doc["metadata"]))

# ------------------------------
# DOCUMENT PROCESSING PIPELINE
# ------------------------------
pipeline = IngestionPipeline(
    transformations=[
        Settings.embed_model  # Only apply the embedding model to pre-chunked text
    ],
)

try:
    # Use the document_objects instead of processed_docs
    nodes = pipeline.run(documents=document_objects)
    
    if not nodes:
        print("❌ No nodes were created. Check document parsing.")
        exit()
    print(f"✅ {len(nodes)} document nodes created and stored in ChromaDB.")

    # Add nodes to ChromaDB
    chroma_path = "./chroma_db"
    chroma_client = PersistentClient(path=chroma_path)
    collection = chroma_client.get_or_create_collection("document_chunks")

    vector_store = ChromaVectorStore(chroma_client, collection_name="document_chunks")
    for idx, node in enumerate(nodes):
    # 1) Sanitize metadata values: convert any Path → str
        clean_meta = {
            key: str(val) if isinstance(val, Path) else val
            for key, val in node.metadata.items()
        }

    # 2) Add to Chroma, using the cleaned metadata
        collection.add(
            ids=[str(idx)],
            documents=[node.text],
            metadatas=[clean_meta]
        )

except Exception as e:
    print(f"❌ Error during ingestion pipeline: {e}")
    exit()

# ------------------------------
# CREATE VECTOR STORE INDEX
# ------------------------------
for node in nodes:
    for k, v in list(node.metadata.items()):
        if isinstance(v, Path):
            node.metadata[k] = str(v)
try:
    index = VectorStoreIndex(nodes, vector_store=vector_store)
    print("✅ Vector store index created successfully.")
    persist_dir = "./persisted_index"
    os.makedirs(persist_dir, exist_ok=True)
    index.storage_context.persist(persist_dir=persist_dir)
    print(f"✅ Index persisted to {persist_dir}")
except Exception as e:
    print(f"❌ Error creating VectorStoreIndex: {e}")
    exit()

# ------------------------------
# CREATE CHAT ENGINE & PROCESS QUERY
# ------------------------------
class RAGQueryEngine(CustomQueryEngine):
    """RAG Query Engine for custom retrieval and response synthesis."""

    retriever: BaseRetriever
    response_synthesizer: BaseSynthesizer

    def custom_query(self, query_str: str):
        nodes = self.retriever.retrieve(query_str)
        response_obj = self.response_synthesizer.synthesize(query_str, nodes)
        return response_obj

retriever = index.as_retriever()
synthesizer = get_response_synthesizer(response_mode="compact")

query_engine = RAGQueryEngine(
    retriever=retriever, response_synthesizer=synthesizer
)

response = query_engine.query("Write a code to find a factorial for a number?")
print(response)
